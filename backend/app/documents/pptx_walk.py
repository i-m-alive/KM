"""Recursive shape walking for PPTX. A plain `for shape in slide.shapes` MISSES
anything inside a grouped shape (MSO_SHAPE_TYPE.GROUP) - PowerPoint decks
routinely group a logo + label, an org chart, or a diagram, and a client name
inside that group was invisible to extraction, masking, AND image scanning
until this existed. All three (extract.py, render.py, images.py) must walk
shapes the same way or detection and masking silently disagree.
"""

from collections.abc import Iterator


def iter_shapes_recursive(shapes) -> Iterator:
    """Yield every shape, descending into group shapes depth-first."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes_recursive(shape.shapes)
