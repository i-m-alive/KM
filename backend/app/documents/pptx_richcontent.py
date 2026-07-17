"""PPTX content that lives outside the shape-tree text walk: chart cache
text, SmartArt diagram text, and embedded/linked OLE object text. Each of
these is a relationship-graph node (chart part, diagram data part, embedded
object part) that `iter_shapes_recursive` walks past without ever descending
into - a client name in a chart's category axis or a SmartArt node was
previously invisible to extraction, detection, masking, AND post-render
verification simultaneously, since all of those reuse the same shape-tree
walk / extract_chunks() call.

Extraction and masking deliberately read/write the SAME xml nodes (see
_chart_cache_text_elements, used by both chart_text_lines and mask_chart) so
what gets proposed for masking is exactly what gets matched later - the two
must never structurally disagree about what "the chart's text" means, the
same invariant pptx_walk.py documents for shape-tree walking.

Callers (extract.py, render.py) inject their own text-matching/masking
callables rather than this module importing render.py's - render.py already
imports this module, so importing back would be circular.
"""

import tempfile
from collections.abc import Callable
from os import fdopen, remove

from lxml import etree

GRAPHIC_DATA_URI_DIAGRAM = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
_DGM_NS = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"

_NUMERIC_ANCESTOR_TAGS = {f"{{{_C_NS}}}numCache", f"{{{_C_NS}}}numLit"}


def is_smartart(shape) -> bool:
    graphicFrame = getattr(shape, "_graphicFrame", None)
    return graphicFrame is not None and graphicFrame.graphicData_uri == GRAPHIC_DATA_URI_DIAGRAM


def is_ole(shape) -> bool:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    return getattr(shape, "shape_type", None) in (
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
        MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
    )


# ---- Charts ----
def _has_numeric_ancestor(el) -> bool:
    for ancestor in el.iterancestors():
        if ancestor.tag in _NUMERIC_ANCESTOR_TAGS:
            return True
    return False


def _chart_cache_text_elements(chart_space) -> list:
    """Every c:v element carrying STRING content - series-name and
    category-label caches (including multi-level categories), wherever they
    sit in the tree. c:v also appears under c:numCache/c:numLit for plotted
    numeric values; those must never be touched by text masking or the
    plotted data itself would be corrupted, so they're the one thing
    excluded here."""
    return [v for v in chart_space.iter(f"{{{_C_NS}}}v") if not _has_numeric_ancestor(v)]


def _chart_formula_refs(chart_space) -> list[str]:
    """c:f cell-reference formula strings (e.g. 'ClientName_Q3!$A$2') -
    surfaced for detection since the reference itself can leak a sheet/tab
    name, but never rewritten during masking: editing a formula reference is
    a fundamentally different, much riskier operation than replacing a
    cached display value, and isn't needed to fix the actual leak (the
    cached c:v text the chart displays is what gets masked)."""
    return sorted({f.text.strip() for f in chart_space.iter(f"{{{_C_NS}}}f") if f.text and f.text.strip()})


def _chart_embedded_sheet_titles(chart) -> list[str]:
    """Sheet/tab names in the chart's embedded worksheet - a real content
    channel (a source tab literally named after the client is a common
    real-world pattern) that the XML cache alone never carries, so without
    this a name visible only via "Edit Data" would never be proposed for
    masking in the first place."""
    try:
        import io

        import openpyxl

        xlsx_part = chart.part.chart_workbook.xlsx_part
        if xlsx_part is None:
            return []
        book = openpyxl.load_workbook(io.BytesIO(xlsx_part.blob))
        return list(book.sheetnames)
    except Exception:
        return []


def chart_text_lines(shape) -> list[str]:
    """Series names, category labels (incl. multi-level), the chart title,
    embedded-worksheet sheet names, and (tagged separately, lower priority)
    formula references."""
    chart = shape.chart
    lines = [v.text.strip() for v in _chart_cache_text_elements(chart._chartSpace) if v.text and v.text.strip()]
    if chart.has_title:
        try:
            title = " ".join(
                run.text for para in chart.chart_title.text_frame.paragraphs for run in para.runs
            ).strip()
            if title:
                lines.append(title)
        except Exception:
            pass
    lines.extend(_chart_embedded_sheet_titles(chart))
    lines.extend(f"[chart cell ref] {f}" for f in _chart_formula_refs(chart._chartSpace))
    return lines


def _mask_chart_embedded_workbook(chart, replace_text_fn: Callable[[str], str]) -> bool:
    """python-pptx/PowerPoint keep the chart's cached display values AND its
    embedded worksheet in sync - masking only the XML cache above would look
    masked in PowerPoint until a user clicks "Edit Data" on the chart, which
    would still show the original value. This can ALSO be the only place a
    surface appears at all (e.g. a sheet/tab name that never made it into
    any cached c:v) - so the caller must call this unconditionally, not only
    when the cache itself changed. Best-effort: if the embedded part is
    missing/corrupt, the visible chart face is still masked, which is a
    strictly smaller exposure than leaving the chart face itself unmasked.
    Returns True if anything in the embedded workbook was changed."""
    try:
        import io

        import openpyxl

        # render.py imports this module, so importing it back at module
        # scope would cycle - deferred import is safe since it only runs
        # once both modules have finished loading.
        from app.documents.render import _safe_sheet_title

        wb_part = chart.part.chart_workbook
        xlsx_part = wb_part.xlsx_part
        if xlsx_part is None:
            return False
        book = openpyxl.load_workbook(io.BytesIO(xlsx_part.blob))
        changed = False
        existing_titles = set(book.sheetnames)
        for ws in book.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str) and cell.value:
                        masked = replace_text_fn(cell.value)
                        if masked != cell.value:
                            cell.value = masked
                            changed = True
            # The sheet/tab NAME itself is a real content channel a cell
            # loop never sees - a chart whose source tab is literally named
            # after the client (a common real-world pattern) would
            # otherwise mask every cell value while leaving the tab name
            # itself visible the moment a user clicks "Edit Data".
            masked_title = replace_text_fn(ws.title)
            if masked_title != ws.title:
                existing_titles.discard(ws.title)
                ws.title = _safe_sheet_title(masked_title, existing_titles)
                existing_titles.add(ws.title)
                changed = True
        if changed:
            buf = io.BytesIO()
            book.save(buf)
            wb_part.update_from_xlsx_blob(buf.getvalue())
        return changed
    except Exception:
        return False


def mask_chart(
    shape,
    replace_text_fn: Callable[[str], str],
    mask_paragraph_fn: Callable[[object], bool],
) -> bool:
    """In-place. `replace_text_fn` mirrors render.py's _replace_text (plain
    string in, masked string out); `mask_paragraph_fn` mirrors
    _mask_paragraph_runs (reused as-is for the chart title, since
    chart.chart_title.text_frame is a normal TextFrame)."""
    chart = shape.chart
    changed = False
    for v in _chart_cache_text_elements(chart._chartSpace):
        if v.text:
            masked = replace_text_fn(v.text)
            if masked != v.text:
                v.text = masked
                changed = True
    if chart.has_title:
        try:
            for para in chart.chart_title.text_frame.paragraphs:
                if mask_paragraph_fn(para):
                    changed = True
        except Exception:
            pass
    # Unconditional: the embedded workbook can carry text (a sheet/tab name,
    # in particular) that never appears in the XML cache at all, so gating
    # this behind the cache having changed would silently skip masking it.
    if _mask_chart_embedded_workbook(chart, replace_text_fn):
        changed = True
    return changed


# ---- SmartArt ----
def smartart_data_part(shape):
    """The related data1.xml Part for a SmartArt graphicFrame, or None."""
    graphicData = shape._graphicFrame.graphicData
    relids = graphicData.find(f"{{{_DGM_NS}}}relIds")
    if relids is None:
        return None
    dm_rid = relids.get(f"{{{_R_NS}}}dm")
    if not dm_rid:
        return None
    try:
        return shape.part.related_part(dm_rid)
    except Exception:
        return None


def smartart_text_lines(shape) -> list[str]:
    part = smartart_data_part(shape)
    if part is None:
        return []
    try:
        root = etree.fromstring(part.blob)
    except Exception:
        return []
    return [t.text.strip() for t in root.iter(f"{{{_A_NS}}}t") if t.text and t.text.strip()]


def mask_smartart(shape, surface_to_token: dict[str, str], style: str) -> bool:
    """In-place. SmartArt text is structurally a paragraph
    (dgm:t/a:p/a:r/a:t) - each a:p's a:t elements are run-boundary-aware
    masked via the SAME shared algorithm ordinary text frames use
    (app.masking.spans.apply_spans_to_runs), so a name split across runs
    still matches AND runs outside the match keep their own content/
    formatting, rather than collapsing the whole match into the first run
    and blanking the rest."""
    from app.masking.spans import apply_spans_to_runs

    part = smartart_data_part(shape)
    if part is None:
        return False
    try:
        root = etree.fromstring(part.blob)
    except Exception:
        return False
    changed = False
    for p_el in root.iter(f"{{{_A_NS}}}p"):
        t_elements = p_el.findall(f"{{{_A_NS}}}r/{{{_A_NS}}}t")
        if apply_spans_to_runs(t_elements, surface_to_token, style):
            changed = True
    if changed:
        part.blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
    return changed


# ---- OLE objects ----
def _guess_ole_extension(prog_id: str | None) -> str | None:
    prog_id = (prog_id or "").lower()
    if "word" in prog_id:
        return ".docx"
    if "excel" in prog_id or "sheet" in prog_id:
        return ".xlsx"
    if "powerpoint" in prog_id:
        return ".pptx"
    return None


def _embedded_office_text_lines(blob: bytes, prog_id: str | None) -> list[str]:
    """Recurse into the standard extraction path for an embedded (not
    linked), zip-based Office document. An OLE Compound-File payload
    (D0CF11E0 magic - old-format .doc/.xls, or a non-Office embed) isn't
    parsed here; the caller still gets the shape's display name."""
    ext = _guess_ole_extension(prog_id)
    if ext is None:
        return []
    from app.documents.extract import extract_chunks

    fd, tmp_path = tempfile.mkstemp(suffix=ext)
    try:
        with fdopen(fd, "wb") as f:
            f.write(blob)
        chunks = extract_chunks(tmp_path, "", tmp_path)
        return [line.strip() for c in chunks for line in c.text.split("\n") if line.strip()]
    except Exception:
        return []
    finally:
        try:
            remove(tmp_path)
        except OSError:
            pass


def ole_text_lines(shape) -> list[str]:
    """The shape's own display name (PowerPoint shows this in the Selection
    Pane / alt-text, and a user renaming an OLE object shape to something
    like 'Bajaj_Financials.xlsx' is a real, observed pattern) plus, for an
    embedded Office document, its own extracted text."""
    lines = []
    name = (shape.name or "").strip()
    if name:
        lines.append(name)

    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
        return lines  # linked object: no local payload to read
    try:
        blob = shape.ole_format.blob
    except Exception:
        blob = None
    if not blob or blob[:4] != b"PK\x03\x04":
        return lines  # OLE Compound-File payload (or unreadable) - name is all we get
    lines.extend(_embedded_office_text_lines(blob, shape.ole_format.prog_id))
    return lines


def _mask_embedded_office_blob(blob: bytes, prog_id: str | None, surface_to_token: dict, style: str) -> bytes | None:
    """Re-render an embedded Office document's own content through the
    SAME per-format renderer used for a top-level upload, so an embedded
    doc gets the identical run-boundary-safe masking a standalone file
    would - not a separate, weaker implementation. Deferred import: render.py
    imports this module, so importing it back at module scope would cycle."""
    ext = _guess_ole_extension(prog_id)
    if ext is None:
        return None
    from app.documents.render import _render_docx, _render_pptx, _render_xlsx

    renderer = {".docx": _render_docx, ".pptx": _render_pptx, ".xlsx": _render_xlsx}[ext]
    fd_src, src_path = tempfile.mkstemp(suffix=ext)
    dst_path = src_path + ".masked" + ext
    try:
        with fdopen(fd_src, "wb") as f:
            f.write(blob)
        renderer(src_path, dst_path, surface_to_token, style)
        with open(dst_path, "rb") as f:
            return f.read()
    except Exception:
        return None
    finally:
        for p in (src_path, dst_path):
            try:
                remove(p)
            except OSError:
                pass


def mask_ole(shape, replace_text_fn: Callable[[str], str], surface_to_token: dict, style: str) -> bool:
    """In-place. Masks the shape's display name always; for an embedded
    zip-based Office document, also replaces the embedded payload with a
    fully re-rendered (masked) copy - otherwise the display name could be
    clean while opening the embedded object still reveals the original
    text."""
    changed = False
    name = shape.name
    if name:
        masked = replace_text_fn(name)
        if masked != name:
            shape.name = masked
            changed = True

    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
        return changed
    try:
        blob = shape.ole_format.blob
    except Exception:
        blob = None
    if not blob or blob[:4] != b"PK\x03\x04":
        return changed
    masked_blob = _mask_embedded_office_blob(blob, shape.ole_format.prog_id, surface_to_token, style)
    if masked_blob is None:
        return changed
    blob_rId = shape._graphicFrame.graphicData.blob_rId
    if not blob_rId:
        return changed
    shape.part.related_part(blob_rId).blob = masked_blob
    return True
