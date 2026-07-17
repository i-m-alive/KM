"""Embedded-image extraction. Text extraction (extract.py) never sees these -
a logo or a screenshot with a client name in it is pixels, not text, and was
invisible to the whole pipeline until this module existed.

Enumeration has two layers for OOXML (docx/pptx/xlsx):
  1. The shape tree (slides/paragraphs/cells) - gives precise, human-readable
     locations ("slide 3").
  2. A raw zip glob over */media/* (+ docProps/thumbnail.*) - catches
     everything the shape tree can't reach: slide masters/layouts, headers/
     footers, SmartArt drawing parts, and orphaned media with no live
     relationship chain. All OOXML binary media lives in a shared media/
     folder by OPC convention regardless of which XML part references it, so
     this glob is a MORE complete inventory than relationship-graph walking,
     not less - a media part with a broken/missing rel is still sitting in
     the zip and is exactly the kind of residual content this exists to catch.
  Layer 2 entries are deduped against layer 1 by partname and get a generic,
  path-derived location label since we don't know which slide/section used them.

Each ImageRef carries enough to re-find the same image deterministically at
apply()-time (locator), so we never need to persist image bytes across the
detect()/apply() boundary - we just re-extract from the untouched original.
"""

import io
import os
import zipfile
from dataclasses import dataclass


@dataclass
class ImageRef:
    index: int  # stable order across the whole document
    location_label: str  # e.g. "slide 3", "page 7", "paragraph area 2"
    image_bytes: bytes
    image_format: str  # "png" | "jpeg" | "emf" | "wmf" | "svg" | ...
    locator: dict  # format-specific: how to find/replace this exact image later


VECTOR_FORMATS = {"emf", "wmf", "emz", "wmz", "svg"}


def guess_image_format(data: bytes, fallback: str | None = "png") -> str | None:
    """Sniff the real format from magic bytes - the only reliable source of
    truth. A container's declared extension (OOXML partname, etc.) can lie
    (e.g. a JPEG saved into a part named *.png after a re-export), so any
    caller that hands these bytes to something format-sensitive (Bedrock
    vision, Pillow) must trust this over a declared/assumed format."""
    if data[:8] == b"\x89PNG\r\n\x1a\n":
        return "png"
    if data[:3] == b"\xff\xd8\xff":
        return "jpeg"
    if data[:6] in (b"GIF87a", b"GIF89a"):
        return "gif"
    if data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return "webp"
    if data[:4] == b"\x01\x00\x00\x00" and data[40:44] == b" EMF":
        return "emf"
    if data[:4] == b"\xd7\xcd\xc6\x9a":
        return "wmf"
    if data[:2] == b"\x1f\x8b":
        return "emz"  # gzip-wrapped EMF/WMF - can't tell which without inflating; soffice handles both by extension
    if data[:2] == b"BM":
        return "bmp"
    if data[:4] in (b"II*\x00", b"MM\x00*"):
        return "tiff"
    stripped = data.lstrip()[:256]
    if stripped[:5] == b"<?xml" or stripped[:4] == b"<svg" or b"<svg" in stripped[:256]:
        return "svg"
    return fallback


def _format_for_partname(partname: str, data: bytes) -> str:
    """OOXML always carries the real extension in the partname - trust that
    first (it's how the package itself typed the part) and only fall back to
    magic-byte sniffing when the extension is missing/ambiguous."""
    ext = os.path.splitext(partname)[1].lstrip(".").lower()
    if ext in {"png", "jpeg", "jpg", "gif", "webp", "emf", "wmf", "emz", "wmz", "svg", "bmp", "tiff"}:
        return "jpeg" if ext == "jpg" else ext
    return guess_image_format(data)


_LABEL_RULES = [
    ("docProps/thumbnail", "document thumbnail"),
    ("slideMasters", "slide master graphic"),
    ("slideLayouts", "slide layout graphic"),
    ("diagrams", "SmartArt graphic"),
    ("header", "header graphic"),
    ("footer", "footer graphic"),
]


def _label_for_partname(partname: str) -> str:
    for needle, label in _LABEL_RULES:
        if needle in partname:
            return label
    return "embedded image (not reachable from body text)"


def _glob_media_parts(path: str, media_prefixes: tuple[str, ...]) -> list[tuple[str, bytes]]:
    """Every file under the given media-folder prefixes, plus docProps/thumbnail.*,
    read directly from the zip - independent of any shape tree or rels graph."""
    found: list[tuple[str, bytes]] = []
    try:
        with zipfile.ZipFile(path) as z:
            for name in z.namelist():
                if name.startswith(media_prefixes) or name.startswith("docProps/thumbnail"):
                    try:
                        found.append((name, z.read(name)))
                    except Exception:
                        continue
    except Exception:
        pass
    return found


def _extract_pdf_images(path: str) -> list[ImageRef]:
    import fitz  # PyMuPDF

    refs: list[ImageRef] = []
    doc = fitz.open(path)
    idx = 0
    for page_number, page in enumerate(doc):
        seen_xrefs = set()
        for img in page.get_images(full=True):
            xref = img[0]
            if xref in seen_xrefs:
                continue
            seen_xrefs.add(xref)
            try:
                base = doc.extract_image(xref)
            except Exception:
                continue
            rects = [tuple(r) for r in page.get_image_rects(xref)]
            refs.append(
                ImageRef(
                    index=idx,
                    location_label=f"page {page_number + 1}",
                    image_bytes=base["image"],
                    image_format=base.get("ext", "png"),
                    locator={"kind": "pdf", "page_number": page_number, "xref": xref, "rects": rects},
                )
            )
            idx += 1
    doc.close()
    return refs


def _extract_docx_images(path: str) -> list[ImageRef]:
    import docx

    document = docx.Document(path)
    refs: list[ImageRef] = []
    seen_partnames: set[str] = set()
    idx = 0
    for shape in document.inline_shapes:
        try:
            rid = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
            image_part = document.part.related_parts[rid]
            data = image_part.blob
            partname = str(image_part.partname).lstrip("/")
        except Exception:
            continue
        seen_partnames.add(partname)
        refs.append(
            ImageRef(
                index=idx,
                location_label=f"inline image {idx + 1}",
                image_bytes=data,
                image_format=_format_for_partname(partname, data),
                locator={"kind": "docx", "partname": partname},
            )
        )
        idx += 1

    # Layer 2: raw media glob - catches floating/anchored images (not inline
    # shapes), header/footer graphics, and the docProps thumbnail.
    for partname, data in _glob_media_parts(path, ("word/media/",)):
        if partname in seen_partnames:
            continue
        seen_partnames.add(partname)
        refs.append(
            ImageRef(
                index=idx,
                location_label=_label_for_partname(partname),
                image_bytes=data,
                image_format=_format_for_partname(partname, data),
                locator={"kind": "docx", "partname": partname},
            )
        )
        idx += 1
    return refs


def _pptx_picture_parts(shapes) -> list:
    """Every (partname, data, blob_getter) for PICTURE shapes in a shape tree,
    walking into groups. Shared by the per-slide walk and the layout/master
    walk below - same extraction logic, different source shape trees."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from app.documents.pptx_walk import iter_shapes_recursive

    found = []
    for shape in iter_shapes_recursive(shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            data = shape.image.blob
            rid = shape._element.blip_rId
            image_part = shape.part.related_part(rid)
            partname = str(image_part.partname).lstrip("/")
        except Exception:
            continue
        found.append((partname, data))
    return found


def _show_master_sp(element) -> bool:
    """PowerPoint's per-slide/per-layout "Hide Background Graphics" toggle.
    Both <p:sld> and <p:sldLayout> carry an optional showMasterSp attribute
    (true when absent) - when a slide (or the layout it's built on) has this
    explicitly set to "0", shapes inherited from the layout/master do NOT
    actually render on that slide, even though python-pptx's shape tree can
    still see them via slide.slide_layout/slide_layout.slide_master. Ignoring
    this flag is exactly what turned a two-slide logo placement into a
    reported occurrence on every slide that merely uses the same layout."""
    return element is not None and element.get("showMasterSp") != "0"


def _extract_pptx_images(path: str) -> list[ImageRef]:
    from pptx import Presentation

    prs = Presentation(path)
    refs: list[ImageRef] = []
    seen_partnames: set[str] = set()
    idx = 0
    for slide_number, slide in enumerate(prs.slides):
        for partname, data in _pptx_picture_parts(slide.shapes):
            seen_partnames.add(partname)
            refs.append(
                ImageRef(
                    index=idx,
                    location_label=f"slide {slide_number + 1}",
                    image_bytes=data,
                    image_format=_format_for_partname(partname, data),
                    locator={"kind": "pptx", "partname": partname, "slide_number": slide_number},
                )
            )
            idx += 1

        # Branding placed on the slide LAYOUT or its MASTER (a logo set once
        # in the template rather than pasted onto every slide individually)
        # renders on every slide that uses that layout, but was previously
        # invisible to this per-slide walk entirely - the shape simply isn't
        # in slide.shapes, it belongs to a different part. The raw media glob
        # below still finds the underlying file, but only as ONE physical
        # part with a generic label ("slide layout graphic"), not as N
        # per-slide locations - wildly undercounting a logo that visually
        # appears on every slide (observed: a top-right corner client logo
        # on all 21 slides counted as only 2 occurrences total). Emitting one
        # ImageRef per slide here - same bytes each time - lets the existing
        # SHA-256 dedup in image_scan.py naturally collapse these into one
        # group whose occurrence count and location list are now accurate.
        #
        # BUT this must not fire for a slide that has "hide background
        # graphics" turned on (showMasterSp="0") - that slide genuinely does
        # not render the layout/master picture, so counting it as an
        # occurrence over-reports both occurrence_count and locations. A
        # slide's own showMasterSp gates BOTH the layout's and the master's
        # graphics (if the slide hides background graphics, neither shows);
        # the layout's own showMasterSp additionally gates whether the
        # MASTER's graphics reach slides using that layout at all.
        layout = slide.slide_layout
        master = layout.slide_master if layout is not None else None
        slide_shows_background = _show_master_sp(slide._element)
        layout_shows_master = _show_master_sp(layout._element) if layout is not None else True
        sources = []
        if slide_shows_background:
            if layout is not None:
                sources.append(layout)
            if master is not None and layout_shows_master:
                sources.append(master)
        for source in sources:
            for partname, data in _pptx_picture_parts(source.shapes):
                seen_partnames.add(partname)
                refs.append(
                    ImageRef(
                        index=idx,
                        location_label=f"slide {slide_number + 1}",
                        image_bytes=data,
                        image_format=_format_for_partname(partname, data),
                        locator={"kind": "pptx", "partname": partname, "slide_number": slide_number},
                    )
                )
                idx += 1

    # Layer 2: raw media glob - SmartArt drawing parts, orphaned media with no
    # live relationship chain, and anything else neither walk above can reach.
    for partname, data in _glob_media_parts(path, ("ppt/media/",)):
        if partname in seen_partnames:
            continue
        seen_partnames.add(partname)
        refs.append(
            ImageRef(
                index=idx,
                location_label=_label_for_partname(partname),
                image_bytes=data,
                image_format=_format_for_partname(partname, data),
                locator={"kind": "pptx", "partname": partname, "slide_number": None},
            )
        )
        idx += 1
    return refs


def _extract_xlsx_images(path: str) -> list[ImageRef]:
    """Raw media glob only, deliberately not openpyxl's ws._images: openpyxl
    silently drops WMF and any PIL-unopenable image on load, so ws._images is
    an incomplete inventory - exactly the kind of gap this module exists to
    close. (Redaction uses a different, position-matched path - see render.py -
    since openpyxl renumbers every image partname on save, so a partname
    captured here would not match post-render.)"""
    refs: list[ImageRef] = []
    for idx, (partname, data) in enumerate(_glob_media_parts(path, ("xl/media/",))):
        refs.append(
            ImageRef(
                index=idx,
                location_label="worksheet image",
                image_bytes=data,
                image_format=_format_for_partname(partname, data),
                locator={"kind": "xlsx", "partname": partname},
            )
        )
    return refs


def extract_images(stored_path: str, content_type: str, filename: str) -> list[ImageRef]:
    """Deterministic, order-stable extraction of every embedded image."""
    lower = filename.lower()
    if content_type == "application/pdf" or lower.endswith(".pdf"):
        return _extract_pdf_images(stored_path)
    if lower.endswith(".docx") or "wordprocessingml" in content_type:
        return _extract_docx_images(stored_path)
    if lower.endswith(".pptx") or "presentationml" in content_type:
        return _extract_pptx_images(stored_path)
    if lower.endswith(".xlsx") or "spreadsheetml" in content_type:
        return _extract_xlsx_images(stored_path)
    return []


def image_dimensions(image_bytes: bytes) -> tuple[int, int]:
    from PIL import Image

    with Image.open(io.BytesIO(image_bytes)) as im:
        return im.size
