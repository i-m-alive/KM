"""Render a masked document in the SAME format the user uploaded.

Given the original file and the approved surface->mask-token map, produce a
sanitized copy: masked PDF for a PDF, masked DOCX for a DOCX, masked PPTX for a
PPTX, masked XLSX for an XLSX. Falls back to a .txt only for unknown types.
"""

import io
import os
import re

from app.config import get_settings
from app.documents.pptx_walk import iter_shapes_recursive
from app.masking.pattern import surface_pattern
from app.masking.style import DEFAULT_MASKING_STYLE, replacement_for

settings = get_settings()


def _replace_text(text: str, surface_to_token: dict[str, str], style: str = DEFAULT_MASKING_STYLE) -> str:
    """Case-insensitive, longest-surface-first replacement (mirrors apply_masks)."""
    if not text:
        return text
    out = text
    for surface in sorted(surface_to_token.keys(), key=len, reverse=True):
        out = re.sub(
            surface_pattern(surface),
            lambda m: replacement_for(m.group(0), surface_to_token[surface], style),
            out,
            flags=re.IGNORECASE,
        )
    return out


def _mask_paragraph_runs(p, surface_to_token: dict[str, str], style: str = DEFAULT_MASKING_STYLE) -> bool:
    """Shared by DOCX and PPTX: rewrite a paragraph's runs so the concatenated
    text is masked. Returns True if anything changed."""
    joined = "".join(run.text for run in p.runs)
    if not joined:
        return False
    masked = _replace_text(joined, surface_to_token, style)
    if masked == joined or not p.runs:
        return False
    # Put the whole masked text in the first run; clear the rest (keeps the
    # paragraph's leading formatting; surfaces often span runs anyway).
    p.runs[0].text = masked
    for run in p.runs[1:]:
        run.text = ""
    return True


# ---- DOCX ----
def _docx_mask_table(table, surface_to_token: dict[str, str], style: str) -> None:
    for row in table.rows:
        for cell in row.cells:
            for p in cell.paragraphs:
                _mask_paragraph_runs(p, surface_to_token, style)
            for nested in cell.tables:  # tables nested inside a cell
                _docx_mask_table(nested, surface_to_token, style)


def _docx_mask_textboxes(root_element, parent, surface_to_token: dict[str, str], style: str) -> None:
    """Mask paragraphs inside text boxes (w:txbxContent) - python-docx's
    .paragraphs never descends into these, so cover-page callouts and shape
    labels were previously rendered UNMASKED even when the same name was
    masked everywhere else. Wrapping each w:p in a Paragraph proxy reuses the
    exact run-joining masking the normal body path gets (a surface split
    across runs still matches)."""
    from docx.oxml.ns import qn
    from docx.text.paragraph import Paragraph

    for txbx in root_element.iter(qn("w:txbxContent")):
        for p_el in txbx.iter(qn("w:p")):
            _mask_paragraph_runs(Paragraph(p_el, parent), surface_to_token, style)


def _scrub_docx_note_parts(path: str, surface_to_token: dict[str, str], style: str) -> None:
    """Footnotes/endnotes: python-docx has no API for these parts at all, so
    they're scrubbed in the saved zip via element-text-only substitution (the
    same mechanism comment_scrub uses, with the same limitation: a surface
    split across two runs inside a footnote isn't fixable here - extraction
    still sees it joined, so verification will flag it rather than pass it)."""
    import re as _re
    import shutil
    import zipfile

    from app.documents.comment_scrub import _scrub_element_text

    targets = (r"^word/footnotes\.xml$", r"^word/endnotes\.xml$")
    tmp_path = path + ".notetmp"
    with zipfile.ZipFile(path, "r") as zin:
        replacements: dict[str, bytes] = {}
        for name in zin.namelist():
            if not any(_re.match(t, name) for t in targets):
                continue
            try:
                text = zin.read(name).decode("utf-8")
            except UnicodeDecodeError:
                continue
            new_text, changed = _scrub_element_text(text, ["w:t", "w:delText"], surface_to_token, style)
            if changed:
                replacements[name] = new_text.encode("utf-8")
        if not replacements:
            return
        with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                zout.writestr(item, replacements.get(item.filename, zin.read(item.filename)))
    shutil.move(tmp_path, path)


def _render_docx(src: str, dst: str, surface_to_token: dict[str, str], style: str) -> None:
    import docx

    document = docx.Document(src)

    for p in document.paragraphs:
        _mask_paragraph_runs(p, surface_to_token, style)
    for table in document.tables:
        _docx_mask_table(table, surface_to_token, style)

    # Text boxes anywhere in the body (cover pages, callouts, shape labels).
    _docx_mask_textboxes(document.element.body, document, surface_to_token, style)

    # Headers/footers - letterhead-style templates often carry the client
    # name here, and this was previously never touched at all.
    for section in document.sections:
        for part in (section.header, section.footer):
            for p in part.paragraphs:
                _mask_paragraph_runs(p, surface_to_token, style)
            for table in part.tables:
                _docx_mask_table(table, surface_to_token, style)
            # Text boxes inside headers/footers too.
            _docx_mask_textboxes(part._element, part, surface_to_token, style)

    document.save(dst)

    # Footnotes/endnotes live in parts python-docx can't reach - scrub the
    # saved zip directly.
    _scrub_docx_note_parts(dst, surface_to_token, style)


# ---- PPTX ----
def _pptx_mask_table(table, surface_to_token: dict[str, str], style: str) -> None:
    for row in table.rows:
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                _mask_paragraph_runs(p, surface_to_token, style)


def _render_pptx(src: str, dst: str, surface_to_token: dict[str, str], style: str) -> None:
    from pptx import Presentation

    prs = Presentation(src)

    for slide in prs.slides:
        # Recurse into grouped shapes - the same gap fixed in extract.py.
        # Without this, anything inside a PowerPoint "Group" (logo+label,
        # org charts, diagrams) is silently left unmasked in the output file
        # even though it may have been detected via the slide's full text.
        for shape in iter_shapes_recursive(slide.shapes):
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    _mask_paragraph_runs(p, surface_to_token, style)
            elif getattr(shape, "has_table", False):
                _pptx_mask_table(shape.table, surface_to_token, style)
        # Speaker notes can also carry client-identifying text.
        if slide.has_notes_slide:
            for p in slide.notes_slide.notes_text_frame.paragraphs:
                _mask_paragraph_runs(p, surface_to_token, style)

    prs.save(dst)


# ---- XLSX ----
_SHEET_NAME_FORBIDDEN = re.compile(r"[\[\]:*?/\\]")


def _safe_sheet_title(masked: str, existing: set[str]) -> str:
    """Excel forbids []:*?/\\ in sheet names (so a raw '[CLIENT_1]' token is
    invalid) and caps them at 31 chars; also must be unique per workbook."""
    title = _SHEET_NAME_FORBIDDEN.sub("", masked).strip() or "Sheet"
    title = title[:31]
    base, n = title, 2
    while title in existing:
        suffix = f" ({n})"
        title = base[: 31 - len(suffix)] + suffix
        n += 1
    return title


def _mask_xlsx_sheet_names_and_headers(wb, surface_to_token: dict[str, str], style: str) -> None:
    """Sheet tab names and print headers/footers are real content channels
    ('Bajaj FY24' as a tab name is a leak) that no cell iteration ever sees.
    NOTE: renaming a sheet does not rewrite formulas that reference it by its
    old name - but those formula strings are themselves masked by the cell
    pass (a formula's text is cell.value on a non-data_only load), so the old
    name doesn't survive there either."""
    existing = set(wb.sheetnames)
    for ws in wb.worksheets:
        masked = _replace_text(ws.title, surface_to_token, style)
        if masked != ws.title:
            existing.discard(ws.title)
            ws.title = _safe_sheet_title(masked, existing)
            existing.add(ws.title)
        for hf in (ws.oddHeader, ws.evenHeader, ws.firstHeader, ws.oddFooter, ws.evenFooter, ws.firstFooter):
            if hf is None:
                continue
            for side in (hf.left, hf.center, hf.right):
                text = getattr(side, "text", None)
                if text:
                    masked_text = _replace_text(text, surface_to_token, style)
                    if masked_text != text:
                        side.text = masked_text


def _render_xlsx(
    src: str, dst: str, surface_to_token: dict[str, str], style: str, approved_image_refs: list | None = None
) -> int:
    """Text masking AND image redaction happen in this ONE load->mutate->save
    pass, deliberately - openpyxl renumbers every image's media partname on
    every save(), so a partname captured from a separate, earlier extraction
    would not match this file after it's been rendered once (the exact
    silent-failure mode this whole feature exists to fix, just relocated to
    xlsx). Matching is by raw byte content instead: ws._images[i]._data()
    returns the same bytes as the source zip's media part for any format
    openpyxl can load, since neither loading nor this in-place mutation
    involves a save cycle beforehand. Returns images_redacted count."""
    import openpyxl

    wb = openpyxl.load_workbook(src)
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and cell.value:
                    masked = _replace_text(cell.value, surface_to_token, style)
                    if masked != cell.value:
                        cell.value = masked
    _mask_xlsx_sheet_names_and_headers(wb, surface_to_token, style)

    images_redacted = 0
    if approved_image_refs:
        from app.documents.image_redact import placeholder_png

        approved_bytes = {ref.image_bytes for ref in approved_image_refs if ref.locator.get("kind") == "xlsx"}
        for ws in wb.worksheets:
            for im in ws._images:
                try:
                    # .getvalue() reads the buffer without consuming/closing it -
                    # Image._data() does close its underlying stream as a side
                    # effect, which would corrupt this image at save time if we
                    # called it here just to compare bytes.
                    data = im.ref.getvalue()
                except Exception:
                    continue
                if data in approved_bytes:
                    im.ref = io.BytesIO(placeholder_png(im.width, im.height))
                    im.format = "png"
                    images_redacted += 1

    wb.save(dst)
    return images_redacted


# ---- PDF (redaction) ----
def _norm_pdf_word(w: str) -> str:
    """Word-comparison key: casefolded, punctuation stripped from both ends
    (a doc word 'BAJAJ,' must match the surface 'BAJAJ')."""
    return re.sub(r"^[\W_]+|[\W_]+$", "", w).casefold()


def _word_token_match(doc_norm: str, token: str) -> bool:
    # Possessives are the one common interior-punctuation case: "BAJAJ's"
    # ends in a word char so end-stripping can't remove it.
    return doc_norm == token or doc_norm == token + "'s" or doc_norm == token + "’s"


def _pdf_surface_sequences(page, surface: str) -> list[list]:
    """True word-level occurrences of `surface` on the page. Replaces the old
    page.search_for() + context-padding heuristic, which had two real failure
    modes: (a) a multi-word name wrapped across a line break never matched at
    all (search_for is a single-line substring search), and (b) in tight
    table/column layouts the padded-context word-boundary check could glue an
    adjacent cell's text onto the match and wrongly REJECT a legitimate
    occurrence - leaving it unmasked while verification (which re-extracts
    text with different word segmentation) didn't always catch it.

    Matching is done on page.get_text("words"): consecutive words within the
    same text block (so unrelated columns are never stitched together) whose
    normalized forms equal the surface's tokens. Word-boundary safety is by
    construction - whole words only, so "RIA" can never match inside
    "MATERIAL". Returns a list of matched sequences, each a list of
    fitz.Rect (one per word, possibly spanning lines)."""
    import fitz

    tokens = [t for t in (_norm_pdf_word(t) for t in surface.split()) if t]
    if not tokens:
        return []
    # (x0, y0, x1, y1, text, block_no, line_no, word_no); punctuation-only
    # words are dropped from matching so "Johnson & Johnson" still matches
    # when '&' is its own word.
    words = [(fitz.Rect(w[:4]), _norm_pdf_word(w[4]), w[5]) for w in page.get_text("words")]
    words = [w for w in words if w[1]]

    sequences: list[list] = []
    n = len(tokens)
    for i in range(len(words) - n + 1):
        window = words[i : i + n]
        if all(_word_token_match(w[1], tok) for w, tok in zip(window, tokens)):
            if len({w[2] for w in window}) == 1:  # same block only
                sequences.append([w[0] for w in window])
    return sequences


def _render_pdf(src: str, dst: str, surface_to_token: dict[str, str], style: str) -> None:
    import fitz  # PyMuPDF

    doc = fitz.open(src)
    for page in doc:
        # Longest first so "Acme Corp" is redacted before "Acme".
        for surface in sorted(surface_to_token.keys(), key=len, reverse=True):
            token = surface_to_token[surface]
            for sequence in _pdf_surface_sequences(page, surface):
                for j, rect in enumerate(sequence):
                    if style == "black":
                        # Solid black redaction bar - no replacement text, nothing readable survives.
                        page.add_redact_annot(rect, fill=(0, 0, 0))
                    elif style == "remove":
                        # Blank the region out entirely - no marker left behind.
                        page.add_redact_annot(rect, fill=(1, 1, 1))
                    elif j == 0:
                        # White-out the original text and overlay the traceable
                        # mask token - on the FIRST word of the sequence only,
                        # so a multi-word name yields one token, not one per word.
                        page.add_redact_annot(rect, text=token, fill=(1, 1, 1), text_color=(0, 0, 0), fontsize=8)
                    else:
                        page.add_redact_annot(rect, fill=(1, 1, 1))
        page.apply_redactions()
    doc.save(dst, garbage=4, deflate=True)
    doc.close()


def render_masked_document(
    run_id: str,
    src_path: str,
    content_type: str,
    filename: str,
    surface_to_token: dict[str, str],
    style: str = DEFAULT_MASKING_STYLE,
    approved_image_refs: list | None = None,
) -> tuple[str, int]:
    """Produce a masked copy in the original format. Returns (path,
    xlsx_images_redacted) - the second element is only ever non-zero for
    xlsx, where image redaction must happen in this same pass (see
    _render_xlsx); every other format redacts images separately, after this
    call, via image_redact.redact_images()."""
    out_dir = os.path.join(settings.OUTPUTS_DIR, "sanitization")
    os.makedirs(out_dir, exist_ok=True)
    stem, ext = os.path.splitext(os.path.basename(filename))
    ext = ext.lower()
    dst = os.path.join(out_dir, f"{run_id}__{stem}.sanitized{ext or '.txt'}")

    xlsx_images_redacted = 0
    lower = filename.lower()
    if content_type == "application/pdf" or lower.endswith(".pdf"):
        _render_pdf(src_path, dst, surface_to_token, style)
    elif lower.endswith(".docx") or "wordprocessingml" in content_type:
        _render_docx(src_path, dst, surface_to_token, style)
    elif lower.endswith(".pptx") or "presentationml" in content_type:
        _render_pptx(src_path, dst, surface_to_token, style)
    elif lower.endswith(".xlsx") or "spreadsheetml" in content_type:
        xlsx_images_redacted = _render_xlsx(src_path, dst, surface_to_token, style, approved_image_refs)
    else:
        with open(dst, "w", encoding="utf-8") as f:
            f.write("\n".join(f"{k} -> {v}" for k, v in surface_to_token.items()))
    return dst, xlsx_images_redacted
