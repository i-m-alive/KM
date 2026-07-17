"""Structural extraction + chunking for uploaded documents.

Chunks preserve structure (page for PDF, paragraph-group for DOCX) so a location
stays meaningful and Sanitization can apply masks by exact character offset.
"""

from dataclasses import dataclass

from app.config import get_settings
from app.documents import pptx_richcontent as rich
from app.documents.images import _show_master_sp
from app.documents.pptx_walk import iter_shapes_recursive

settings = get_settings()


@dataclass
class Chunk:
    chunk_id: int
    kind: str  # "page" | "paragraph_group" | "slide" | "sheet"
    label: str  # human-readable, e.g. "page 3"
    text: str


class UnsupportedDocumentError(Exception):
    pass


def _extract_pdf(path: str) -> list[Chunk]:
    import pdfplumber

    chunks: list[Chunk] = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            chunks.append(Chunk(chunk_id=i, kind="page", label=f"page {i + 1}", text=text))
    return chunks


def _pptx_table_lines(table) -> list[str]:
    lines = []
    for row in table.rows:
        cells = [c.text.strip() for c in row.cells if c.text.strip()]
        if cells:
            lines.append(" | ".join(cells))
    return lines


def _pptx_shape_content_lines(shapes) -> tuple[list[str], list[str], list[str], list[str]]:
    """Every text-bearing kind of content in a shape tree - text frames,
    tables, charts, SmartArt, and OLE objects - as four separate line lists
    (text, chart, smartart, ole). Shared by the per-slide walk and the
    layout/master walk so both see exactly the same shape kinds - before
    this, the layout/master walk only handled text frames/tables, so a
    chart/SmartArt/OLE object placed directly on a layout or master (rare,
    but the same asymmetry images.py originally had for pictures) was
    invisible to extraction even though render.py's masking pass (which
    reuses _mask_pptx_shapes for both slide and layout/master shape trees)
    still walks those same shapes looking for something to mask."""
    text_lines: list[str] = []
    chart_lines: list[str] = []
    smartart_lines: list[str] = []
    ole_lines: list[str] = []
    # Recurse into grouped shapes - a plain top-level loop misses any text
    # nested inside a PowerPoint "Group" (very common for logo+label
    # graphics, org charts, diagrams).
    for shape in iter_shapes_recursive(shapes):
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    text_lines.append(text)
        elif getattr(shape, "has_table", False):
            text_lines.extend(_pptx_table_lines(shape.table))
        elif getattr(shape, "has_chart", False):
            chart_lines.extend(rich.chart_text_lines(shape))
        elif rich.is_smartart(shape):
            smartart_lines.extend(rich.smartart_text_lines(shape))
        elif rich.is_ole(shape):
            ole_lines.extend(rich.ole_text_lines(shape))
    return text_lines, chart_lines, smartart_lines, ole_lines


def _extract_pptx(path: str) -> list[Chunk]:
    from pptx import Presentation

    prs = Presentation(path)
    chunks: list[Chunk] = []
    notes_lines: list[str] = []
    # Text placeholders inherited from the slide LAYOUT or MASTER (a footer,
    # a confidentiality line, a template-level client name set once) - the
    # same asymmetry images.py already fixed for pictures (_show_master_sp),
    # applied to text: a slide's OWN shape tree never includes shapes
    # inherited from its layout/master, so this text was previously
    # invisible to extraction (and therefore detection, masking, and
    # verification) entirely. Emitted once per slide that actually shows it
    # (not deduped to once per distinct layout), so occurrence counts stay
    # accurate rather than under-reporting a name shown on every slide as a
    # single hit - the same reasoning images.py's per-slide ImageRef
    # emission already documents.
    layout_lines: list[str] = []
    master_lines: list[str] = []
    # Charts, SmartArt, and OLE objects are relationship-graph nodes (a chart
    # part, a diagram data part, an embedded object part) that the shape-tree
    # walk below never descends into on its own - without pptx_richcontent,
    # a client name in a chart's category axis or a SmartArt node was
    # invisible to extraction, detection, masking, AND post-render
    # verification simultaneously, since all of those reuse this function.
    # Collected as their own channels (same convention as speaker notes
    # below), tagged per-slide, rather than merged into each slide's own text.
    chart_lines: list[str] = []
    smartart_lines: list[str] = []
    ole_lines: list[str] = []
    # A layout/master's OWN content never varies per slide - only the
    # per-slide LABEL does - so it's computed once per distinct
    # layout/master (keyed by identity) and reused for every slide that
    # shares it, rather than re-walking the same shape tree for every one
    # of potentially many slides built on the same handful of layouts.
    layout_content_cache: dict[int, tuple[list[str], list[str], list[str], list[str]]] = {}
    master_content_cache: dict[int, tuple[list[str], list[str], list[str], list[str]]] = {}

    for i, slide in enumerate(prs.slides):
        text_lines, slide_chart, slide_smartart, slide_ole = _pptx_shape_content_lines(slide.shapes)
        chart_lines.extend(f"[slide {i + 1} chart] {t}" for t in slide_chart)
        smartart_lines.extend(f"[slide {i + 1} smartart] {t}" for t in slide_smartart)
        ole_lines.extend(f"[slide {i + 1} embedded object] {t}" for t in slide_ole)
        # Preserve slide boundaries so Deck Drafting (A-06) can reuse them later.
        chunks.append(Chunk(chunk_id=i, kind="slide", label=f"slide {i + 1}", text="\n".join(text_lines)))

        # Layout/master content - gated by the SAME showMasterSp flag
        # images.py already established: a slide's own showMasterSp gates
        # BOTH the layout's and the master's inherited graphics (if the
        # slide hides background graphics, neither shows), and the
        # layout's own showMasterSp additionally gates whether the
        # MASTER's graphics reach slides using that layout at all.
        layout = slide.slide_layout
        master = layout.slide_master if layout is not None else None
        slide_shows_background = _show_master_sp(slide._element)
        layout_shows_master = _show_master_sp(layout._element) if layout is not None else True
        if slide_shows_background and layout is not None:
            key = id(layout)
            if key not in layout_content_cache:
                layout_content_cache[key] = _pptx_shape_content_lines(layout.shapes)
            l_text, l_chart, l_smartart, l_ole = layout_content_cache[key]
            layout_lines.extend(f"[slide {i + 1} layout] {t}" for t in l_text)
            chart_lines.extend(f"[slide {i + 1} layout chart] {t}" for t in l_chart)
            smartart_lines.extend(f"[slide {i + 1} layout smartart] {t}" for t in l_smartart)
            ole_lines.extend(f"[slide {i + 1} layout embedded object] {t}" for t in l_ole)
            if master is not None and layout_shows_master:
                key = id(master)
                if key not in master_content_cache:
                    master_content_cache[key] = _pptx_shape_content_lines(master.shapes)
                m_text, m_chart, m_smartart, m_ole = master_content_cache[key]
                master_lines.extend(f"[slide {i + 1} master] {t}" for t in m_text)
                chart_lines.extend(f"[slide {i + 1} master chart] {t}" for t in m_chart)
                smartart_lines.extend(f"[slide {i + 1} master smartart] {t}" for t in m_smartart)
                ole_lines.extend(f"[slide {i + 1} master embedded object] {t}" for t in m_ole)

        # Speaker notes are a real content channel presenters use for client
        # context - render.py already masks them, but extraction (and
        # therefore NER pre-pass, the LLM Detector, the dictionary full-text
        # sweep, AND post-render verification, which all reuse this function)
        # previously never read them at all. A client name that ONLY appears
        # in a note was invisible end-to-end: never proposed for masking, and
        # verification couldn't have caught the miss either since it can't
        # see text it never extracts. Appended as one extra chunk (not merged
        # into the slide's own chunk_id) so existing per-slide numbering -
        # which Deck Drafting is expected to rely on - is untouched.
        if slide.has_notes_slide:
            note_text = "\n".join(
                p.text.strip() for p in slide.notes_slide.notes_text_frame.paragraphs if p.text.strip()
            )
            if note_text:
                notes_lines.append(f"[slide {i + 1} notes] {note_text}")
    if chart_lines:
        chunks.append(Chunk(chunk_id=len(chunks), kind="chart", label="chart text", text="\n".join(chart_lines)))
    if smartart_lines:
        chunks.append(Chunk(chunk_id=len(chunks), kind="smartart", label="smartart text", text="\n".join(smartart_lines)))
    if ole_lines:
        chunks.append(Chunk(chunk_id=len(chunks), kind="ole", label="embedded objects", text="\n".join(ole_lines)))
    if layout_lines:
        chunks.append(Chunk(chunk_id=len(chunks), kind="layout-text", label="layout text", text="\n".join(layout_lines)))
    if master_lines:
        chunks.append(Chunk(chunk_id=len(chunks), kind="master-text", label="master text", text="\n".join(master_lines)))
    if notes_lines:
        chunks.append(Chunk(chunk_id=len(chunks), kind="notes", label="speaker notes", text="\n".join(notes_lines)))
    return chunks


def _docx_table_lines(table) -> list[str]:
    lines = []
    for row in table.rows:
        for cell in row.cells:
            for p in cell.paragraphs:
                if p.text.strip():
                    lines.append(p.text.strip())
            # Tables nested inside a cell (common in complex templates).
            for nested in cell.tables:
                lines.extend(_docx_table_lines(nested))
    return lines


def _docx_textbox_lines(root_element) -> list[str]:
    """Text inside text boxes (w:txbxContent). python-docx's .paragraphs never
    descends into these - and Word cover pages / callouts put client names in
    text boxes constantly - so without this walk the text is invisible to
    detection, masking, AND verification: a silent leak, not a flagged one."""
    from docx.oxml.ns import qn

    lines = []
    for txbx in root_element.iter(qn("w:txbxContent")):
        for p in txbx.iter(qn("w:p")):
            text = "".join(t.text or "" for t in p.iter(qn("w:t"))).strip()
            if text:
                lines.append(text)
    return lines


def _docx_part_text(path: str, partnames: tuple[str, ...]) -> list[str]:
    """All w:t text from raw zip parts python-docx has no API for
    (footnotes, endnotes)."""
    import re as _re
    import xml.etree.ElementTree as ET
    import zipfile

    W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    lines: list[str] = []
    try:
        with zipfile.ZipFile(path) as z:
            names = z.namelist()
            for pattern in partnames:
                for name in names:
                    if not _re.match(pattern, name):
                        continue
                    try:
                        root = ET.fromstring(z.read(name))
                    except ET.ParseError:
                        continue
                    for p in root.iter(f"{W}p"):
                        text = "".join(t.text or "" for t in p.iter(f"{W}t")).strip()
                        if text:
                            lines.append(text)
    except Exception:
        pass
    return lines


def _extract_docx(path: str) -> list[Chunk]:
    import docx

    document = docx.Document(path)
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]

    # Headers/footers carry client names surprisingly often (letterhead-style
    # templates) and were previously invisible to the whole pipeline.
    header_footer_lines: list[str] = []
    for section in document.sections:
        for part in (section.header, section.footer):
            for p in part.paragraphs:
                if p.text.strip():
                    header_footer_lines.append(p.text.strip())

    # Text boxes anywhere in the body, plus footnotes/endnotes - all outside
    # python-docx's .paragraphs and previously invisible end to end.
    textbox_lines = _docx_textbox_lines(document.element.body)
    note_lines = _docx_part_text(path, (r"^word/footnotes\.xml$", r"^word/endnotes\.xml$"))

    # Tables at the document body level (and any tables nested within them).
    table_lines: list[str] = []
    for table in document.tables:
        table_lines.extend(_docx_table_lines(table))

    chunks: list[Chunk] = []
    if header_footer_lines:
        chunks.append(Chunk(chunk_id=-1, kind="header_footer", label="headers & footers", text="\n".join(header_footer_lines)))
    if textbox_lines:
        chunks.append(Chunk(chunk_id=-2, kind="text_boxes", label="text boxes", text="\n".join(textbox_lines)))
    if note_lines:
        chunks.append(Chunk(chunk_id=-3, kind="notes", label="footnotes & endnotes", text="\n".join(note_lines)))

    # Group body paragraphs into chunks so each LLM call sees a coherent block.
    per_chunk = max(1, settings.CHUNK_PARAGRAPHS)
    for i in range(0, len(paragraphs), per_chunk):
        group = paragraphs[i : i + per_chunk]
        start = i // per_chunk
        chunks.append(
            Chunk(
                chunk_id=start,
                kind="paragraph_group",
                label=f"paragraphs {i + 1}-{i + len(group)}",
                text="\n\n".join(group),
            )
        )

    if table_lines:
        chunks.append(Chunk(chunk_id=len(chunks), kind="tables", label="tables", text="\n".join(table_lines)))

    return chunks


def _xlsx_header_footer_lines(ws) -> list[str]:
    """Print header/footer text - set from Page Layout in Excel, carries
    'Client X - Confidential'-style lines surprisingly often, and never
    appears in any cell."""
    lines = []
    for hf in (ws.oddHeader, ws.evenHeader, ws.firstHeader, ws.oddFooter, ws.evenFooter, ws.firstFooter):
        if hf is None:
            continue
        for side in (hf.left, hf.center, hf.right):
            text = (getattr(side, "text", None) or "").strip()
            if text:
                lines.append(text)
    return lines


def _extract_xlsx(path: str) -> list[Chunk]:
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True)
    chunks: list[Chunk] = []
    for i, sheet_name in enumerate(wb.sheetnames):
        ws = wb[sheet_name]
        # The sheet TAB NAME itself is content ("Bajaj FY24" as a tab name is
        # a leak) - include it in the chunk text so detection, masking, and
        # verification all see it, not just the human-readable label.
        lines: list[str] = [sheet_name]
        lines.extend(_xlsx_header_footer_lines(ws))
        for row in ws.iter_rows():
            cells = [str(c.value).strip() for c in row if c.value is not None and str(c.value).strip()]
            if cells:
                lines.append(" | ".join(cells))
        chunks.append(Chunk(chunk_id=i, kind="sheet", label=f"sheet '{sheet_name}'", text="\n".join(lines)))
    return chunks


def extract_chunks(stored_path: str, content_type: str, filename: str) -> list[Chunk]:
    """Extract a document into structural chunks. Dispatches by content type / extension."""
    lower = filename.lower()
    if content_type == "application/pdf" or lower.endswith(".pdf"):
        return _extract_pdf(stored_path)
    if (
        content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        or lower.endswith(".docx")
    ):
        return _extract_docx(stored_path)
    if (
        content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or lower.endswith(".pptx")
    ):
        return _extract_pptx(stored_path)
    if (
        content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        or lower.endswith(".xlsx")
    ):
        return _extract_xlsx(stored_path)
    raise UnsupportedDocumentError(
        f"Unsupported document type '{content_type}' ({filename}). Sanitization supports PDF, DOCX, PPTX, and XLSX."
    )
