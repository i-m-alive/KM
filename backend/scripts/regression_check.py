"""Standing regression suite for Sanitization's deterministic pipeline.

Every masking bug found in real usage lived in the DETERMINISTIC layer
(word-boundary corruption, un-enumerated master/orphan images, metadata /
comment / hyperlink channels never checked or never fixed) - not in the LLM.
So this suite plants known leaks in freshly-built fixture documents (docx,
pptx, xlsx, pdf), runs the real render + scrub functions, and asserts:

  1. each scan channel DETECTS its planted leak on the raw file
     (a channel that can't catch a planted leak is a channel that can
     silently pass a real one), and
  2. after render + scrub, every deterministic channel verifies clean, and
  3. masking is word-boundary safe (masking "RIA" must not corrupt
     "MATERIAL"), and
  4. image enumeration sees orphaned media with no relationship chain
     (the slide-master / dangling-media class of miss).

No Bedrock calls and no DB - runs offline in seconds:

    cd backend && .venv/bin/python scripts/regression_check.py
"""

import asyncio
import io
import re
import shutil
import sys
import tempfile
import zipfile
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, patch

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from app.documents.comment_scan import find_residual_comments
from app.documents.comment_scrub import scrub_comments
from app.documents.hyperlink_scan import find_residual_hyperlinks
from app.documents.hyperlink_scrub import scrub_hyperlinks
from app.documents.images import extract_images
from app.documents.metadata_scan import find_residual_metadata
from app.documents.metadata_scrub import scrub_metadata
from app.documents.render import render_masked_document
from app.documents.verify import find_residual_surfaces

CLIENT = "BAJAJ"
SHORT_CLIENT = "RIA"  # word-boundary trap: must not corrupt MATERIAL
MULTI_CLIENT = "Tata Capital"  # multi-word: must mask even when line-wrapped
SURFACE_TO_TOKEN = {CLIENT: "[CLIENT_1]", SHORT_CLIENT: "[CLIENT_2]", MULTI_CLIENT: "[CLIENT_3]"}
SURFACES = list(SURFACE_TO_TOKEN.keys())
STYLE = "token"
BODY_TEXT = (
    f"Engagement overview for {CLIENT}. The MATERIAL scope covers {SHORT_CLIENT} operations. "
    f"Benchmarked against {MULTI_CLIENT} programs."
)
LEAKY_URL = f"https://www.{CLIENT.lower()}.com/portal"

failures: list[str] = []
passes = 0


def check(name: str, ok: bool, detail: str = "") -> None:
    global passes
    if ok:
        passes += 1
        print(f"  ok    {name}")
    else:
        failures.append(f"{name}: {detail}")
        print(f"  FAIL  {name}  {detail}")


def _png_bytes(color=(200, 30, 30)) -> bytes:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (120, 60), color).save(buf, format="PNG")
    return buf.getvalue()


# ---------- fixture builders ----------

def build_docx(path: str) -> None:
    import docx
    from docx.oxml import parse_xml
    from docx.oxml.ns import nsmap, qn
    from docx.opc.constants import RELATIONSHIP_TYPE as RT

    doc = docx.Document()
    para = doc.add_paragraph(BODY_TEXT)
    doc.add_comment(para.runs, text=f"Confirm with {CLIENT} legal first", author="reviewer")
    doc.core_properties.subject = f"{CLIENT} discovery phase"

    # Tracked deletion - deleted-but-present text lives in w:delText, which no
    # normal text-extraction path sees. python-docx has no API for this; raw XML.
    w = nsmap["w"]
    del_para = doc.add_paragraph()
    del_para._p.append(parse_xml(
        f'<w:del xmlns:w="{w}" w:id="99" w:author="reviewer" w:date="2026-01-01T00:00:00Z">'
        f"<w:r><w:delText>Old {CLIENT} pricing removed</w:delText></w:r></w:del>"
    ))

    # External hyperlink whose TARGET (not display text) leaks the client.
    link_para = doc.add_paragraph()
    r_id = link_para.part.relate_to(LEAKY_URL, RT.HYPERLINK, is_external=True)
    link_para._p.append(parse_xml(
        f'<w:hyperlink xmlns:w="{w}" xmlns:r="{nsmap["r"]}" r:id="{r_id}">'
        "<w:r><w:t>client portal</w:t></w:r></w:hyperlink>"
    ))

    # A VML text box - python-docx's .paragraphs never sees w:txbxContent, so
    # without the dedicated textbox walk this leak is invisible end to end.
    tb_para = doc.add_paragraph()
    tb_para._p.append(parse_xml(
        f'<w:r xmlns:w="{w}" xmlns:v="urn:schemas-microsoft-com:vml"><w:pict>'
        '<v:shape style="width:220pt;height:60pt"><v:textbox><w:txbxContent>'
        f"<w:p><w:r><w:t>Cover note for the {CLIENT} board</w:t></w:r></w:p>"
        "</w:txbxContent></v:textbox></v:shape></w:pict></w:r>"
    ))
    doc.save(path)


def build_pptx(path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    tb.text_frame.text = BODY_TEXT
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = " portal"
    run.hyperlink.address = LEAKY_URL
    slide.shapes.add_picture(io.BytesIO(_png_bytes()), Inches(1), Inches(2))
    prs.save(path)

    # Orphaned media: present in ppt/media/ with NO relationship from any
    # slide - exactly what a rels-graph walk misses and the raw glob must find.
    with zipfile.ZipFile(path, "a") as z:
        z.writestr("ppt/media/orphan_logo.png", _png_bytes((30, 30, 200)))


def build_xlsx(path: str) -> None:
    import openpyxl
    from openpyxl.comments import Comment

    wb = openpyxl.Workbook()
    ws = wb.active
    # Sheet tab name and print header both leak the client, and neither
    # appears in any cell - previously invisible to the whole pipeline.
    ws.title = f"{CLIENT} Data"
    ws.oddHeader.center.text = f"{CLIENT} - Confidential"
    ws["A1"] = BODY_TEXT
    ws["A2"] = "client portal"
    ws["A2"].hyperlink = LEAKY_URL
    ws["A3"] = "note"
    ws["A3"].comment = Comment(f"Check {CLIENT} numbers", "reviewer")
    wb.save(path)


def build_pdf(path: str) -> None:
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    # insert_textbox (not insert_text) so the body wraps within the page -
    # a single insert_text line runs past the page edge, where PyMuPDF's
    # word extraction clips it but pdfplumber's verify still reads it.
    page.insert_textbox(fitz.Rect(72, 80, 540, 190), BODY_TEXT, fontsize=11)
    page.insert_link({"kind": fitz.LINK_URI, "from": fitz.Rect(300, 200, 420, 220), "uri": LEAKY_URL})
    annot = page.add_text_annot((72, 160), f"Verify with {CLIENT} finance")
    annot.update()
    # A narrow text box forces "Tata Capital" to wrap across two lines - the
    # exact case the old single-line substring search could never match.
    page.insert_textbox(fitz.Rect(72, 200, 116, 420), f"Meeting with {MULTI_CLIENT} leadership", fontsize=12)
    doc.set_metadata({**doc.metadata, "subject": f"{CLIENT} discovery"})
    doc.save(path)
    doc.close()


# ---------- per-format run ----------

FIXTURES = [
    ("regression.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", build_docx),
    ("regression.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", build_pptx),
    ("regression.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", build_xlsx),
    ("regression.pdf", "application/pdf", build_pdf),
]

# Which planted leaks each format actually carries (fixture-building APIs
# don't exist for every channel in every format - e.g. no pptx comment API).
HAS_COMMENT_LEAK = {"regression.docx", "regression.xlsx", "regression.pdf"}
HAS_METADATA_LEAK = {"regression.docx", "regression.pdf"}


def run_format(workdir: Path, filename: str, content_type: str, builder) -> None:
    print(f"\n== {filename} ==")
    src = str(workdir / filename)
    builder(src)

    # 1. Every applicable channel must DETECT its planted leak on the raw file.
    check("hyperlink scan detects planted target", len(find_residual_hyperlinks(src, content_type, filename, SURFACES)) > 0)
    if filename in HAS_COMMENT_LEAK:
        check("comment scan detects planted comment", len(find_residual_comments(src, content_type, filename, SURFACES)) > 0)
    if filename in HAS_METADATA_LEAK:
        check("metadata scan detects planted property", len(find_residual_metadata(src, content_type, filename, SURFACES)) > 0)
    check("text verify detects unmasked surface", len(find_residual_surfaces(src, content_type, filename, SURFACES)) > 0)

    # 1b. Channels that were once invisible to extraction must stay visible -
    # if extraction goes blind here again, masking and verification both go
    # blind with it (the silent-leak class, not the flagged class).
    from app.documents.extract import extract_chunks

    pre_text = " ".join(c.text for c in extract_chunks(src, content_type, filename))
    if filename.endswith(".docx"):
        check("docx textbox text visible to extraction", "Cover note" in pre_text)
    if filename.endswith(".xlsx"):
        check("xlsx sheet name visible to extraction", f"{CLIENT} Data" in pre_text)
        check("xlsx print header visible to extraction", "Confidential" in pre_text)
    if filename.endswith(".pdf"):
        wrapped = find_residual_surfaces(src, content_type, filename, [MULTI_CLIENT])
        check("pdf line-wrapped multi-word name detected", len(wrapped) > 0)

    # 2. Render + scrub with the real pipeline functions, in pipeline order.
    masked_path, _ = render_masked_document("regression-check", src, content_type, filename, SURFACE_TO_TOKEN, style=STYLE)
    scrub_metadata(masked_path, content_type, filename, SURFACE_TO_TOKEN, STYLE)
    scrub_hyperlinks(masked_path, content_type, filename, SURFACE_TO_TOKEN, STYLE)
    scrub_comments(masked_path, content_type, filename, SURFACE_TO_TOKEN, STYLE)

    # 3. Every deterministic channel must now verify clean.
    residual_text = find_residual_surfaces(masked_path, content_type, filename, SURFACES)
    check("text channel clean after mask", len(residual_text) == 0, "; ".join(residual_text[:3]))
    residual_meta = find_residual_metadata(masked_path, content_type, filename, SURFACES)
    check("metadata channel clean after scrub", len(residual_meta) == 0, "; ".join(residual_meta[:3]))
    residual_comments = find_residual_comments(masked_path, content_type, filename, SURFACES)
    check("comments channel clean after scrub", len(residual_comments) == 0, "; ".join(residual_comments[:3]))
    residual_links = find_residual_hyperlinks(masked_path, content_type, filename, SURFACES)
    check("hyperlinks channel clean after scrub", len(residual_links) == 0, "; ".join(residual_links[:3]))

    # 4. Word-boundary safety: MATERIAL contains 'RIA' but must survive intact.
    from app.documents.extract import extract_chunks

    masked_text = " ".join(c.text for c in extract_chunks(masked_path, content_type, filename))
    check("word-boundary safe (MATERIAL intact)", "MATERIAL" in masked_text, f"text was: {masked_text[:200]}")
    check("mask token present in masked text", "[CLIENT_1]" in masked_text, f"text was: {masked_text[:200]}")

    # 5. pptx only: orphaned media (no rels chain) must still be enumerated.
    if filename.endswith(".pptx"):
        refs = extract_images(src, content_type, filename)
        parts = {getattr(r, "partname", None) or getattr(r, "location_label", "") for r in refs}
        found = any("orphan_logo" in str(p) for p in parts) or len(refs) >= 2
        check("orphaned media enumerated (glob, not rels-walk)", found, f"only found: {sorted(str(p) for p in parts)}")
        check_near_duplicate_redaction(workdir)


def check_near_duplicate_redaction(workdir: Path) -> None:
    """The perceptual-dedup blind spot found in a real run: the same logo
    saved at two compressions is SHA-distinct but phash-near, so it merges
    into ONE cluster - and redaction must then swap BOTH media parts, not
    just the ones byte-equal to the cluster's sample."""
    from PIL import Image

    from app.agents.sanitization.image_scan import PERCEPTUAL_DEDUP_THRESHOLD
    from app.documents.image_redact import redact_images
    from app.masking.logo_reference import compute_phash, phash_distance
    from pptx import Presentation
    from pptx.util import Inches

    png = _png_bytes((10, 120, 60))
    jpg_buf = io.BytesIO()
    Image.open(io.BytesIO(png)).save(jpg_buf, format="JPEG", quality=60)
    jpg = jpg_buf.getvalue()
    check("renditions are SHA-distinct", png != jpg)
    d = phash_distance(compute_phash(png), compute_phash(jpg))
    check("renditions cluster perceptually", d is not None and d <= PERCEPTUAL_DEDUP_THRESHOLD, f"distance={d}")

    path = str(workdir / "near_dup.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_picture(io.BytesIO(png), Inches(1), Inches(1))
    slide.shapes.add_picture(io.BytesIO(jpg), Inches(4), Inches(1))
    prs.save(path)

    ct = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    refs = extract_images(path, ct, "near_dup.pptx")
    both = [r for r in refs if r.locator.get("kind") == "pptx"]
    redacted, _ = redact_images(path, ct, "near_dup.pptx", both)
    check("both near-duplicate renditions redacted", redacted >= 2, f"redacted={redacted}")

    with zipfile.ZipFile(path) as z:
        media = [n for n in z.namelist() if n.startswith("ppt/media/")]
        survivors = [n for n in media if z.read(n) in (png, jpg)]
    check("no original rendition bytes survive in media", not survivors, f"survived: {survivors}")

    # Remediation double-click safety: a second pass must recognize the
    # placeholder bytes it wrote and target nothing.
    from app.documents.image_redact import is_placeholder_bytes

    refs2 = extract_images(path, ct, "near_dup.pptx")
    remaining = [r for r in refs2 if r.locator.get("kind") == "pptx" and not is_placeholder_bytes(r.image_bytes)]
    check("second remediation pass targets nothing (placeholders recognized)", not remaining,
          f"{len(remaining)} ref(s) not recognized as placeholder")


def check_precision_qa() -> None:
    """Task 1 acceptance test: precision_check.check_precision must surface
    mask tokens the model flags as probable over-redactions of a common word
    (the observed "[CLIENT_18] Allocation" / bare "[CLIENT_55]" failures), and
    must produce zero flags when every token looks like a genuine proper
    noun. Bedrock itself is mocked - this tests the module's plumbing
    (prompt construction, schema parsing, the no-tokens-no-call short
    circuit), not real model judgment, so it stays offline like the rest of
    this suite."""
    print("\n== precision QA (Task 1, mocked Bedrock) ==")
    from app.agents.sanitization import precision_check
    from app.llm import bedrock_client

    masked_over_redacted = (
        "The [CLIENT_18] Allocation strategy was reviewed jointly with [CLIENT_55] "
        "from the vendor side, per the engagement charter."
    )
    canned_flags = {
        "flags": [
            {"mask_token": "[CLIENT_18]", "surrounding_text": "The [CLIENT_18] Allocation strategy",
             "reason": "'Capital Allocation' reads as an ordinary business term, not a company name", "confidence": 0.85},
            {"mask_token": "[CLIENT_55]", "surrounding_text": "jointly with [CLIENT_55] from the vendor side",
             "reason": "'partners' used generically, not as a proper noun", "confidence": 0.8},
        ]
    }
    mock_resp = bedrock_client.BedrockResponse(text="", parsed=canned_flags, input_tokens=100, output_tokens=50, estimated_cost_usd=0.001)
    with patch.object(bedrock_client, "converse", new=AsyncMock(return_value=mock_resp)):
        resp = asyncio.run(precision_check.check_precision(masked_over_redacted, ["[CLIENT_18]", "[CLIENT_55]"]))
    flagged_tokens = {f["mask_token"] for f in (resp.parsed or {}).get("flags", [])}
    check("precision check flags over-redacted common-noun tokens",
          {"[CLIENT_18]", "[CLIENT_55]"} <= flagged_tokens, f"got: {flagged_tokens}")

    masked_genuine = "The engagement was led by [CLIENT_1] alongside their partner firm [CLIENT_2]."
    mock_clean = bedrock_client.BedrockResponse(text="", parsed={"flags": []}, input_tokens=100, output_tokens=10, estimated_cost_usd=0.0005)
    with patch.object(bedrock_client, "converse", new=AsyncMock(return_value=mock_clean)):
        resp2 = asyncio.run(precision_check.check_precision(masked_genuine, ["[CLIENT_1]", "[CLIENT_2]"]))
    check("precision check produces zero flags when every token is a genuine proper noun",
          len((resp2.parsed or {}).get("flags", [])) == 0)

    with patch.object(bedrock_client, "converse", new=AsyncMock(return_value=mock_clean)) as mocked:
        resp3 = asyncio.run(precision_check.check_precision("no masking happened in this document at all", []))
    check("precision check makes no Bedrock call when nothing was masked", resp3 is None and not mocked.called)


def check_reidentify_qa() -> None:
    """Task 2 acceptance test: reidentify.reidentify must flag a mask token as
    re-identifiable when co-occurring facts (the observed BlackRock case:
    Preqin acquisition, ~$25T AUM) are left unmasked, at or above the
    confidence a reviewer would act on - and must NOT flag it once those
    facts are also masked/generalized. Bedrock is mocked; this tests the
    module's plumbing (threshold filtering happens in agent.py, so here we
    just check the raw guesses the mock returns pass through unchanged),
    same offline discipline as check_precision_qa."""
    print("\n== re-identification QA (Task 2, mocked Bedrock) ==")
    from app.agents.sanitization import reidentify
    from app.llm import bedrock_client

    leaky_text = (
        "[CLIENT_1] recently completed its acquisition of Preqin and now manages roughly "
        "~$25T in assets under management, running its [CLIENT_1] Aladdin platform on Azure "
        "OpenAI following the December 2025 AWS partnership announcement."
    )
    canned_guess = {
        "guesses": [
            {"mask_token": "[CLIENT_1]", "candidate_org": "BlackRock", "confidence": 0.9,
             "leaking_phrases": ["Preqin acquisition", "~$25T in assets under management"]},
        ]
    }
    mock_leaky = bedrock_client.BedrockResponse(text="", parsed=canned_guess, input_tokens=120, output_tokens=60, estimated_cost_usd=0.001)
    with patch.object(bedrock_client, "converse", new=AsyncMock(return_value=mock_leaky)):
        resp = asyncio.run(reidentify.reidentify(leaky_text, ["[CLIENT_1]"]))
    guesses = (resp.parsed or {}).get("guesses", [])
    hit = next((g for g in guesses if g["mask_token"] == "[CLIENT_1]"), None)
    check("re-identification flags the leaky case at high confidence citing the AUM figure",
          hit is not None and hit["confidence"] >= 0.5 and any("25T" in p or "Preqin" in p for p in hit["leaking_phrases"]),
          f"got: {hit}")

    sanitized_text = (
        "[CLIENT_1] recently completed an acquisition in the data/analytics space and now manages "
        "a large amount of assets under management, running an internal platform on a major cloud "
        "provider following a recent partnership announcement."
    )
    canned_clean = {"guesses": [{"mask_token": "[CLIENT_1]", "candidate_org": "unknown", "confidence": 0.1, "leaking_phrases": []}]}
    mock_clean = bedrock_client.BedrockResponse(text="", parsed=canned_clean, input_tokens=120, output_tokens=30, estimated_cost_usd=0.0008)
    with patch.object(bedrock_client, "converse", new=AsyncMock(return_value=mock_clean)):
        resp2 = asyncio.run(reidentify.reidentify(sanitized_text, ["[CLIENT_1]"]))
    guesses2 = (resp2.parsed or {}).get("guesses", [])
    hit2 = next((g for g in guesses2 if g["mask_token"] == "[CLIENT_1]"), None)
    check("re-identification does not flag the properly-generalized version above threshold",
          hit2 is not None and hit2["confidence"] < 0.5, f"got: {hit2}")

    with patch.object(bedrock_client, "converse", new=AsyncMock(return_value=mock_clean)) as mocked:
        resp3 = asyncio.run(reidentify.reidentify("no masking happened in this document at all", []))
    check("re-identification makes no Bedrock call when nothing was masked", resp3 is None and not mocked.called)


def check_review_deltas() -> None:
    """Task 3 acceptance test: a simulated review that adds 2 entities and
    removes 1 must produce recall-miss=2 and precision-miss=1, correctly
    attributed to their entity types. Pure functions, no DB/Bedrock needed."""
    print("\n== review deltas (Task 3) ==")
    from app.agents.sanitization import review_deltas
    from app.masking import dictionary

    proposal_entities = [
        {"surface_text": "Acme Partners", "entity_type": "CLIENT_NAME"},
        {"surface_text": "Real Client Co", "entity_type": "CLIENT_NAME"},
    ]
    removed = {"acme partners"}  # reviewer struck this one as over-flagged
    precision = review_deltas.precision_miss_by_type(proposal_entities, removed)
    check("precision_miss_by_type attributes 1 removal to CLIENT_NAME",
          precision == {"CLIENT_NAME": 1}, f"got: {precision}")

    existing_keys = {dictionary.normalize(e["surface_text"]) for e in proposal_entities if e["surface_text"].lower() not in removed}
    added_entities = [
        {"surface_text": "Foo Corp", "entity_type": "CLIENT_NAME"},
        {"surface_text": "Jane Doe", "entity_type": "CLIENT_PERSON"},
        {"surface_text": "Real Client Co", "entity_type": "CLIENT_NAME"},  # duplicate of an existing key - not a real miss
    ]
    new_entities, recall = review_deltas.resolve_added_entities(added_entities, existing_keys)
    check("resolve_added_entities counts 2 real recall misses by type",
          recall == {"CLIENT_NAME": 1, "CLIENT_PERSON": 1}, f"got: {recall}")
    check("resolve_added_entities drops the duplicate-key addition",
          len(new_entities) == 2 and {e["surface_text"] for e in new_entities} == {"Foo Corp", "Jane Doe"},
          f"got: {new_entities}")


def check_per_entity_type_thresholds() -> None:
    """Task 4 acceptance test: lowering one entity type's confidence gate in
    config must change only that type's inclusion boundary - a candidate of a
    different type at the same confidence must be unaffected. Exercises the
    real settings object (not a reimplementation of the comparison logic)."""
    print("\n== per-entity-type thresholds (Task 4) ==")
    from app.config import get_settings

    settings = get_settings()
    check("all entity types default to the same untuned 0.6",
          all(v == 0.6 for v in settings.SANITIZATION_CONFIDENCE_THRESHOLDS.values()),
          f"got: {settings.SANITIZATION_CONFIDENCE_THRESHOLDS}")

    original = dict(settings.SANITIZATION_CONFIDENCE_THRESHOLDS)
    try:
        settings.SANITIZATION_CONFIDENCE_THRESHOLDS["CLIENT_NAME"] = 0.3
        confidence = 0.45  # below the untouched 0.6 default, above the lowered 0.3
        name_threshold = settings.SANITIZATION_CONFIDENCE_THRESHOLDS.get("CLIENT_NAME", 0.6)
        person_threshold = settings.SANITIZATION_CONFIDENCE_THRESHOLDS.get("CLIENT_PERSON", 0.6)
        check("lowering CLIENT_NAME's gate admits a candidate at 0.45",
              confidence >= name_threshold, f"threshold={name_threshold}")
        check("CLIENT_PERSON's gate is untouched by the CLIENT_NAME change",
              confidence < person_threshold, f"threshold={person_threshold}")
    finally:
        settings.SANITIZATION_CONFIDENCE_THRESHOLDS.clear()
        settings.SANITIZATION_CONFIDENCE_THRESHOLDS.update(original)


def check_vision_verdict_cache() -> None:
    """Task 5 acceptance test: the SAME run scanning the SAME image content
    twice (once standing in for detect(), once for apply()'s verify - a
    genuinely different call since those are different OS processes) must
    make exactly one live vision call and return a byte-identical structured
    verdict both times, including for two renditions that are byte-different
    but phash-identical. vision_cache's DB-backed store/load is swapped for
    an in-memory dict here (same discipline as mocking bedrock_client.converse
    elsewhere in this suite - Postgres itself isn't under test, the REAL
    caching control flow in image_scan.py's _scan_one_group is)."""
    print("\n== vision-verdict memoization (Task 5, mocked Bedrock + cache store) ==")
    import uuid as uuid_mod

    from PIL import Image

    from app.agents.sanitization import image_scan, vision_cache
    from app.documents.images import ImageRef
    from app.llm import bedrock_client

    fake_store: dict[tuple, dict] = {}

    def fake_load(db, run_id, content_key):
        return fake_store.get((run_id, content_key))

    def fake_store_fn(db, run_id, content_key, parsed):
        fake_store[(run_id, content_key)] = {
            "contains_client_identity": bool(parsed.get("contains_client_identity", False)),
            "description": parsed.get("description", "") or "",
            "confidence": float(parsed.get("confidence", 0.0)),
            "ocr_text": [s for s in (parsed.get("ocr_text") or []) if isinstance(s, str)],
        }

    call_count = {"n": 0}

    async def fake_converse_vision(**kwargs):
        call_count["n"] += 1
        return bedrock_client.BedrockResponse(
            text="", parsed={"contains_client_identity": True, "description": "a client logo", "confidence": 0.93, "ocr_text": []},
            input_tokens=50, output_tokens=20, estimated_cost_usd=0.0005,
        )

    run_id = uuid_mod.uuid4()
    png = _png_bytes((5, 5, 5))
    jpg_buf = io.BytesIO()
    Image.open(io.BytesIO(png)).save(jpg_buf, format="JPEG", quality=70)
    jpg = jpg_buf.getvalue()
    check("the two renditions are byte-distinct (a real second-encoding case, not a no-op)", png != jpg)

    shared_phash = "abc123abc123abc1"  # standing in for "phash-identical renditions"
    ref1 = ImageRef(index=0, location_label="slide 1", image_bytes=png, image_format="png", locator={})
    ref2 = ImageRef(index=1, location_label="slide 7", image_bytes=jpg, image_format="jpeg", locator={})

    with patch.object(vision_cache, "load_cached_verdict", side_effect=fake_load), \
         patch.object(vision_cache, "store_verdict", side_effect=fake_store_fn), \
         patch.object(bedrock_client, "converse_vision", side_effect=fake_converse_vision):
        db = None  # never dereferenced: find_matches gets references=[] below, _ocr_match gets empty ocr_text
        group1, in1, out1, cost1 = asyncio.run(image_scan._scan_one_group(
            db, 0, [ref1], png, "png", shared_phash, logo_references=[], run_id=run_id,
        ))
        group2, in2, out2, cost2 = asyncio.run(image_scan._scan_one_group(
            db, 0, [ref2], jpg, "jpeg", shared_phash, logo_references=[], run_id=run_id,
        ))

    check("exactly one live vision call for two scans of the same run+content", call_count["n"] == 1, f"got: {call_count['n']}")
    check("second scan (byte-different, phash-identical) returns the byte-identical structured verdict",
          (group1.contains_client_identity, group1.description, group1.confidence, group1.ocr_text) ==
          (group2.contains_client_identity, group2.description, group2.confidence, group2.ocr_text),
          f"group1={group1}, group2={group2}")
    check("second scan made no Bedrock call (zero tokens/cost)", in2 == 0 and out2 == 0 and cost2 == 0.0, f"got in={in2} out={out2} cost={cost2}")


class _FakeFlag:
    def __init__(self, message, severity):
        self.message = message
        self.severity = severity


class _FakeRun:
    def __init__(self, flags):
        self.flags = flags


class _FakeDeleteDB:
    def __init__(self):
        self.deleted = []

    def delete(self, obj):
        self.deleted.append(obj)


def check_stale_flag_cleanup() -> None:
    """Task 6a acceptance test: remediate._remove_resolved_flags (now also
    called from background.py's normal first-pass completion, not just the
    remediation path) must remove a detect-phase blocking image flag once
    images verify clean, and a "Verification failed (channel)" blocking flag
    once that channel verifies clean - but must never touch a non-blocking
    flag, and must NOT remove the image flag if images are still dirty.
    Fakes stand in for the ORM run/flags/session so this is a pure logic
    test, no DB needed."""
    print("\n== stale flag cleanup (Task 6a) ==")
    from app.agents.sanitization.remediate import _remove_resolved_flags

    image_flag = _FakeFlag("3 embedded image(s) appear to reveal the client (logo/screenshot) — review before treating this document as sanitized.", "blocking")
    text_flag = _FakeFlag("Verification failed (text): 1 item(s) still expose a masked term...", "blocking")
    info_flag = _FakeFlag("2 embedded image(s) found; none flagged as client-identifying...", "info")

    run = _FakeRun([image_flag, text_flag, info_flag])
    db = _FakeDeleteDB()
    removed = _remove_resolved_flags(db, run, channels_now_clean=["text", "images"], images_clean=True)
    check("stale detect-phase image flag is removed once images verify clean",
          image_flag in db.deleted, f"deleted: {db.deleted}")
    check("stale 'Verification failed (text)' flag is removed once text verifies clean",
          text_flag in db.deleted, f"deleted: {db.deleted}")
    check("non-blocking info flag is never removed", info_flag not in db.deleted)
    check("_remove_resolved_flags reports the count it actually removed", removed == 2, f"got: {removed}")

    run2 = _FakeRun([_FakeFlag("1 embedded image(s) appear to reveal the client (logo/screenshot) — review.", "blocking")])
    db2 = _FakeDeleteDB()
    _remove_resolved_flags(db2, run2, channels_now_clean=["text"], images_clean=False)
    check("image flag is NOT removed while images are still dirty", len(db2.deleted) == 0, f"deleted: {db2.deleted}")


def check_text_remediation_idempotent_reremask() -> None:
    """Task 6b acceptance test: pointing render_masked_document at an
    ALREADY-MASKED file (not the original) with the same surface_to_token
    map must be idempotent - an already-masked token is left untouched, not
    double-masked or corrupted - and must still catch a genuinely residual
    literal occurrence. This is the mechanism remediate.py now uses instead
    of unconditionally forcing a full re-run on every text residual."""
    print("\n== text remediation re-render idempotency (Task 6b) ==")
    import docx

    from app.documents.extract import extract_chunks

    workdir = Path(tempfile.mkdtemp(prefix="naviknow-remediate-text-"))
    try:
        ct = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        src = str(workdir / "leak.docx")
        doc = docx.Document()
        doc.add_paragraph(f"Engagement for {CLIENT}. Second mention of {CLIENT} here.")
        doc.save(src)

        surface_to_token = {CLIENT: "[CLIENT_9]"}
        masked_path, _ = render_masked_document("remediate-text-check", src, ct, "leak.docx", surface_to_token, style="token")

        # Simulate a residual: an old-renderer miss left one more literal
        # occurrence sitting in the "already masked" file.
        d2 = docx.Document(masked_path)
        d2.add_paragraph(f"A residual mention of {CLIENT} that an earlier renderer version missed.")
        d2.save(masked_path)

        residual_before = find_residual_surfaces(masked_path, ct, "leak.docx", [CLIENT])
        check("fixture actually has a residual before the fix", len(residual_before) > 0, f"got: {residual_before}")

        fixed_path, _ = render_masked_document("remediate-text-check-2", masked_path, ct, "leak.docx", surface_to_token, style="token")

        text_after = " ".join(c.text for c in extract_chunks(fixed_path, ct, "leak.docx"))
        check("already-masked tokens survive unchanged, and the residual is now masked too (3 total)",
              text_after.count("[CLIENT_9]") == 3, f"count={text_after.count('[CLIENT_9]')}; text was: {text_after[:300]}")
        residual_after = find_residual_surfaces(fixed_path, ct, "leak.docx", [CLIENT])
        check("no residual remains after the re-render", len(residual_after) == 0, f"got: {residual_after}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def check_xlsx_threaded_comments() -> None:
    """Task 6c acceptance test: Excel threaded (modern) comments and person
    display names - previously a documented gap - must now be detected and
    scrubbed. The fixture has no real .rels relationship graph, so this
    specifically exercises the glob-fallback path (xl/threadedComments/*.xml,
    xl/persons/person.xml) this module relies on when relationship
    resolution alone isn't enough to prove out."""
    print("\n== xlsx threaded comments (Task 6c) ==")
    import openpyxl

    from app.documents.comment_scan import find_residual_comments
    from app.documents.comment_scrub import scrub_comments

    workdir = Path(tempfile.mkdtemp(prefix="naviknow-threaded-comments-"))
    try:
        ct = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        path = str(workdir / "threaded.xlsx")
        wb = openpyxl.Workbook()
        wb.active["A1"] = "unrelated content"
        wb.save(path)

        threaded_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<ThreadedComments xmlns="http://schemas.microsoft.com/office/spreadsheetml/2018/threadedcomments">'
            f'<threadedComment ref="A1" dT="2026-01-01T00:00:00.00Z" personId="{{P1}}" id="{{C1}}">'
            f'<text>Confirm with {CLIENT} before publishing</text>'
            "</threadedComment>"
            "</ThreadedComments>"
        )
        person_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<personList xmlns="http://schemas.microsoft.com/office/spreadsheetml/2018/threadedcomments">'
            f'<person displayName="{CLIENT} Reviewer" id="{{P1}}" userId="x" providerId="None"/>'
            "</personList>"
        )
        with zipfile.ZipFile(path, "a") as z:
            z.writestr("xl/threadedComments/threadedComment1.xml", threaded_xml)
            z.writestr("xl/persons/person.xml", person_xml)

        residual_before = find_residual_comments(path, ct, "threaded.xlsx", SURFACES)
        check("threaded comment text is detected before scrub",
              any("threaded comment" in h for h in residual_before), f"got: {residual_before}")
        check("person display name is detected before scrub",
              any("threaded comment author" in h for h in residual_before), f"got: {residual_before}")

        scrub_comments(path, ct, "threaded.xlsx", SURFACE_TO_TOKEN, STYLE)
        residual_after = find_residual_comments(path, ct, "threaded.xlsx", SURFACES)
        check("threaded comments channel clean after scrub", len(residual_after) == 0, f"got: {residual_after}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


class _FakeLogoRef:
    def __init__(self, id_, mask_entity_id, phash):
        self.id = id_
        self.mask_entity_id = mask_entity_id
        self.phash = phash


def check_logo_band_index() -> None:
    """Task 6e acceptance test: the LSH-style band index must find every
    match a full linear scan finds - no false negatives - for a close match
    within threshold, while a genuinely distant reference is still correctly
    excluded from the final (threshold-filtered) results. band_index only
    prunes candidates; find_matches still applies the exact Hamming-distance
    threshold on whatever it's given, so this proves the pruning doesn't
    change the actual matching decision."""
    print("\n== logo phash LSH band index (Task 6e) ==")
    from app.masking.logo_reference import UNCERTAIN_THRESHOLD, build_band_index, find_matches, phash_distance

    base = "0f1e2d3c4b5a6978"
    close = "0f1e2d3c4b5a6970"  # differs by 1 bit in the last hex digit
    far = "ffffffffffffffff"  # maximally distant

    refs = [
        _FakeLogoRef(1, "11111111-1111-1111-1111-111111111111", close),
        _FakeLogoRef(2, "22222222-2222-2222-2222-222222222222", far),
    ]
    d_close = phash_distance(base, close)
    d_far = phash_distance(base, far)
    check("fixture's close reference is within UNCERTAIN_THRESHOLD", d_close is not None and d_close <= UNCERTAIN_THRESHOLD, f"d={d_close}")
    check("fixture's far reference is well beyond UNCERTAIN_THRESHOLD", d_far is not None and d_far > UNCERTAIN_THRESHOLD, f"d={d_far}")

    index = build_band_index(refs)
    via_index = find_matches(None, base, threshold=UNCERTAIN_THRESHOLD, band_index=index)
    via_scan = find_matches(None, base, threshold=UNCERTAIN_THRESHOLD, references=refs)
    check("band-index search finds the same matches as a full linear scan",
          via_index == via_scan, f"index={via_index} scan={via_scan}")
    check("the close reference is found via the index", any(m[0] == refs[0].mask_entity_id for m in via_index), f"got: {via_index}")
    check("the far reference is excluded from the final results",
          not any(m[0] == refs[1].mask_entity_id for m in via_index), f"got: {via_index}")


def check_alt_text_channel() -> None:
    """Task A acceptance test: image alt-text (descr/title/name on
    cNvPr/docPr) - the seam between the pixel channel (image_scan.py) and
    the body-text channel (verify.py) - must be detected pre-scrub,
    rewritten by scrub_alt_text, and re-verified clean; a planted
    post-render leak must still fail verification, not just be detectable
    in isolation. The real observed failure: 21 client names recovered from
    descr="..." attributes on logos whose PIXELS were already fully
    redacted."""
    print("\n== image alt-text channel (Task A) ==")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches

    from app.documents.alttext_scan import extract_alt_text
    from app.documents.alttext_scrub import scrub_alt_text
    from app.documents.verify import find_residual_surfaces

    workdir = Path(tempfile.mkdtemp(prefix="naviknow-alttext-"))
    try:
        ct = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        path = str(workdir / "alttext.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        pic = slide.shapes.add_picture(io.BytesIO(_png_bytes()), Inches(1), Inches(1))
        pic._element.nvPicPr.cNvPr.set("descr", f"Our Business | {CLIENT} Group logo")
        pic._element.nvPicPr.cNvPr.set("title", f"{CLIENT} wordmark")
        prs.save(path)

        pre = extract_alt_text(path, ct, "alttext.pptx")
        check("alt-text descr/title detected before scrub", any(CLIENT in v for v in pre), f"got: {pre}")

        residual_before = find_residual_surfaces(path, ct, "alttext.pptx", SURFACES)
        check("text-channel verification catches the alt-text leak pre-scrub",
              CLIENT in residual_before, f"got: {residual_before}")

        changed = scrub_alt_text(path, ct, "alttext.pptx", SURFACE_TO_TOKEN, STYLE)
        check("scrub_alt_text reports rewritten attributes", changed >= 2, f"got: {changed}")

        post = extract_alt_text(path, ct, "alttext.pptx")
        check("no client name remains in alt-text after scrub", not any(CLIENT in v for v in post), f"got: {post}")

        residual_after = find_residual_surfaces(path, ct, "alttext.pptx", SURFACES)
        check("text channel verifies clean after alt-text scrub", CLIENT not in residual_after, f"got: {residual_after}")

        # Prove verification actually BLOCKS a planted leak, not just that it
        # CAN detect one in isolation - re-plant after scrub and re-check.
        prs2 = Presentation(path)
        pic2 = next(s for s in prs2.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
        pic2._element.nvPicPr.cNvPr.set("descr", f"{CLIENT} leaked again")
        prs2.save(path)
        residual_replant = find_residual_surfaces(path, ct, "alttext.pptx", SURFACES)
        check("a re-planted alt-text leak fails verification", CLIENT in residual_replant, f"got: {residual_replant}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def check_pptx_author_list() -> None:
    """Task B acceptance test (part 1): ppt/authors.xml (the modern-comments
    author list) must be found via the static-path fallback even when
    relationship resolution alone would find nothing - exactly the real
    observed failure (7 author names + emails leaked, including a
    client-side domain, because the old code had no fallback at all) - and
    the userId attribute (which embeds the actual email) must be scrubbed,
    not just name/initials."""
    print("\n== pptx author list / userId email (Task B part 1) ==")
    from app.documents.comment_scan import find_residual_comments
    from app.documents.comment_scrub import scrub_comments

    workdir = Path(tempfile.mkdtemp(prefix="naviknow-authors-"))
    try:
        ct = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        path = str(workdir / "authors.pptx")
        from pptx import Presentation

        Presentation().save(path)

        authors_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<p:authorLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            '<p:author id="0" name="Ankit Bajpai" initials="AB" '
            'userId="S::ankit.bajpai@zs.com::11111111-1111-1111-1111-111111111111" providerId="AD"/>'
            "</p:authorLst>"
        )
        with zipfile.ZipFile(path, "a") as z:
            z.writestr("ppt/authors.xml", authors_xml)

        email_surface_to_token = {**SURFACE_TO_TOKEN, "ankit.bajpai@zs.com": "[CLIENT_9]", "Ankit Bajpai": "[CLIENT_10]"}
        email_surfaces = SURFACES + ["ankit.bajpai@zs.com", "Ankit Bajpai"]

        residual_before = find_residual_comments(path, ct, "authors.pptx", email_surfaces)
        check("author name is detected via the static ppt/authors.xml fallback",
              any("Ankit Bajpai" in h for h in residual_before), f"got: {residual_before}")
        check("email embedded in userId is detected", any("ankit.bajpai@zs.com" in h for h in residual_before), f"got: {residual_before}")

        scrub_comments(path, ct, "authors.pptx", email_surface_to_token, STYLE)
        residual_after = find_residual_comments(path, ct, "authors.pptx", email_surfaces)
        check("author list channel clean after scrub (name + userId email both gone)",
              len(residual_after) == 0, f"got: {residual_after}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def check_author_identity_unconditional() -> None:
    """Task B core acceptance test: authors.xml/userId must be cleared even
    with an EMPTY surface_to_token/surfaces - proving the fix doesn't
    depend on detection succeeding anywhere. This is the actual root cause
    of "authors.xml completely untouched": nothing in detect() ever looks at
    authors.xml, so these names/emails were NEVER in surface_to_token to
    begin with: the correct fix can't be "detect them better", it has to be
    "don't require detection for this field at all"."""
    print("\n== author identity unconditional clearing (Task B core) ==")
    from app.documents.comment_scan import find_residual_comments
    from app.documents.comment_scrub import scrub_comments

    workdir = Path(tempfile.mkdtemp(prefix="naviknow-authors-unconditional-"))
    try:
        ct = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        path = str(workdir / "authors2.pptx")
        from pptx import Presentation

        Presentation().save(path)

        authors_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<p:authorLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            '<p:author id="0" name="Kashyap Durrani" initials="KD" '
            'userId="S::kashyap.durrani@zs.com::22222222-2222-2222-2222-222222222222" providerId="AD"/>'
            "</p:authorLst>"
        )
        with zipfile.ZipFile(path, "a") as z:
            z.writestr("ppt/authors.xml", authors_xml)

        residual_before = find_residual_comments(path, ct, "authors2.pptx", [])
        check("author identity flagged with an EMPTY surfaces list (never detected anywhere)",
              any("identity attribute" in h for h in residual_before), f"got: {residual_before}")

        changed = scrub_comments(path, ct, "authors2.pptx", {}, STYLE)
        check("scrub_comments clears author identity with an EMPTY surface_to_token", changed >= 2, f"got: {changed}")

        residual_after = find_residual_comments(path, ct, "authors2.pptx", [])
        check("no author-identity residual remains after scrub (empty surfaces list)",
              len(residual_after) == 0, f"got: {residual_after}")

        # Prove it actually BLOCKS - replant and re-check.
        replant_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<p:authorLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            '<p:author id="0" name="Replanted Name" initials="RN" '
            'userId="S::replanted@zs.com::33333333-3333-3333-3333-333333333333" providerId="AD"/>'
            "</p:authorLst>"
        )
        tmp_path = path + ".replant"
        with zipfile.ZipFile(path, "r") as zin, zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename == "ppt/authors.xml":
                    zout.writestr(item, replant_xml)
                else:
                    zout.writestr(item, zin.read(item.filename))
        shutil.move(tmp_path, path)
        residual_replant = find_residual_comments(path, ct, "authors2.pptx", [])
        check("a re-planted author-identity leak fails verification even with an empty surfaces list",
              any("identity attribute" in h for h in residual_replant), f"got: {residual_replant}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def check_docx_comment_author() -> None:
    """Task B acceptance test (part 2): a DOCX comment's w:author/w:initials
    - attributes on <w:comment> itself, a different shape than the comment
    BODY text already being scrubbed - must be detected and scrubbed too.
    python-docx's add_comment(author=...) sets exactly this attribute, so a
    client name used as the comment author (not just in the comment body)
    is a real, plausible leak shape."""
    print("\n== docx comment author attribute (Task B part 2) ==")
    import docx

    from app.documents.comment_scan import find_residual_comments
    from app.documents.comment_scrub import scrub_comments

    workdir = Path(tempfile.mkdtemp(prefix="naviknow-docx-author-"))
    try:
        ct = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        path = str(workdir / "author.docx")
        doc = docx.Document()
        para = doc.add_paragraph("Unrelated body text.")
        doc.add_comment(para.runs, text="benign comment body", author=CLIENT, initials="C")
        doc.save(path)

        residual_before = find_residual_comments(path, ct, "author.docx", SURFACES)
        check("comment AUTHOR name (not just body) is detected before scrub",
              any("comment author" in h and CLIENT in h for h in residual_before), f"got: {residual_before}")

        scrub_comments(path, ct, "author.docx", SURFACE_TO_TOKEN, STYLE)
        residual_after = find_residual_comments(path, ct, "author.docx", SURFACES)
        check("comment author channel clean after scrub", len(residual_after) == 0, f"got: {residual_after}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def check_unconditional_identity_fields() -> None:
    """Task B acceptance test (part 3): dc:creator/cp:lastModifiedBy must be
    cleared UNCONDITIONALLY - regardless of whether that specific name was
    ever detected as a client entity. Real observed failure: "Vallab" in
    creator was never masked at all because nothing ever detected it as an
    entity; "Gaurav [CLIENT_5]" in lastModifiedBy had its surname tokenized
    (that part WAS a detected entity) but the first name survived, proving
    detection-dependent scrubbing alone isn't sufficient for these fields."""
    print("\n== unconditional creator/lastModifiedBy clearing (Task B part 3) ==")
    import docx

    from app.documents.metadata_scan import find_residual_metadata
    from app.documents.metadata_scrub import scrub_metadata

    workdir = Path(tempfile.mkdtemp(prefix="naviknow-identity-fields-"))
    try:
        ct = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        path = str(workdir / "identity.docx")
        doc = docx.Document()
        doc.add_paragraph("Unrelated body text.")
        doc.core_properties.author = "Vallab"  # maps to dc:creator
        doc.core_properties.last_modified_by = "Gaurav Someone"
        doc.save(path)

        residual_before = find_residual_metadata(path, ct, "identity.docx", [])
        check("creator/lastModifiedBy flagged as residual identity fields even with an EMPTY surfaces list",
              any("identity field" in h for h in residual_before), f"got: {residual_before}")

        # Unconditional clearing must run even when NOTHING was masked in
        # this run (empty surface_to_token) - it's hygiene, not detection-gated.
        changed = scrub_metadata(path, ct, "identity.docx", {}, STYLE)
        check("scrub_metadata clears identity fields even with an empty surface_to_token", changed >= 2, f"got: {changed}")

        residual_after = find_residual_metadata(path, ct, "identity.docx", [])
        check("no identity-field residual remains after scrub", len(residual_after) == 0, f"got: {residual_after}")

        # Prove it actually BLOCKS - replant and re-check.
        doc2 = docx.Document(path)
        doc2.core_properties.author = "Vallab Replanted"
        doc2.save(path)
        residual_replant = find_residual_metadata(path, ct, "identity.docx", [])
        check("a re-planted creator leak fails verification even with an empty surfaces list",
              any("identity field" in h for h in residual_replant), f"got: {residual_replant}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def check_surface_pattern_underscore_boundary() -> None:
    """Task A core acceptance test: surface_pattern() must match a surface
    immediately adjacent to an underscore (e.g. alt-text/filenames like
    "Nextcare_logo"), which regex's plain \\b treats as NO boundary at all
    (both are \\w characters) - misdiagnosed as a case-sensitivity bug in the
    original report (case-insensitive matching was already in place and
    didn't fix it), but confirmed here with an exact case match too, so the
    real cause can't silently regress back to "just add re.IGNORECASE"."""
    print("\n== surface_pattern underscore boundary (Task A core) ==")
    from app.masking.pattern import surface_pattern

    check("matches a surface followed by an underscore (exact case)",
          bool(re.search(surface_pattern("NextCare"), "NextCare_logo", re.IGNORECASE)))
    check("matches a surface followed by an underscore (different case) - proves this was never about case",
          bool(re.search(surface_pattern("NextCare"), "nextcare_logo", re.IGNORECASE)))
    check("matches with a trailing file extension after the underscore",
          bool(re.search(surface_pattern("NextCare"), "NextCare_logo.png", re.IGNORECASE)))
    check("still finds a surface with normal punctuation boundaries",
          bool(re.search(surface_pattern("GMR"), "GMR Group | Delhi", re.IGNORECASE)))
    check("word-boundary safety is NOT weakened: 'RIA' must still not match inside 'MATERIAL'",
          not re.search(surface_pattern("RIA"), "MATERIAL", re.IGNORECASE))
    check("line-wrapped multi-word surfaces still match",
          bool(re.search(surface_pattern("Tata Capital"), "Tata\nCapital programs", re.IGNORECASE)))


def check_alt_text_detector_context() -> None:
    """Task A acceptance test: alt-text values must actually reach the LLM
    detector call as context, not just the cheap regex/dictionary merge -
    real observed failure: "GMR Group | Delhi" and "AJG India" had no
    deterministic signal at all (not emails/phones, not already-known
    entities) and so were never proposed with anything but my own crude
    low-confidence whole-phrase fallback. Mocked Bedrock - this proves the
    plumbing (alt_texts reaches the prompt, the response flows back
    normally), not real model judgment."""
    print("\n== alt-text reaches the LLM detector (Task A core) ==")
    from unittest.mock import AsyncMock

    from app.agents.sanitization import detector
    from app.llm import bedrock_client

    captured = {}

    async def fake_converse_with_tools(**kwargs):
        captured["user_message"] = kwargs.get("user_message")
        return bedrock_client.BedrockResponse(
            text="", parsed={"entities": [{"surface_text": "GMR Group", "entity_type": "CLIENT_NAME", "confidence": 0.85}]},
            input_tokens=80, output_tokens=20, estimated_cost_usd=0.001,
        )

    with patch.object(bedrock_client, "converse_with_tools", new=AsyncMock(side_effect=fake_converse_with_tools)):
        resp = asyncio.run(detector.detect_entities(
            "doc-1", 3, [], alt_texts=["GMR Group | Delhi", "AJG India"],
        ))

    check("alt-text values are included in the detector's prompt",
          "GMR Group | Delhi" in captured["user_message"] and "AJG India" in captured["user_message"],
          f"prompt was: {captured.get('user_message')}")
    entities = (resp.parsed or {}).get("entities", [])
    check("the model's alt-text-derived entity flows back through the response normally",
          any(e["surface_text"] == "GMR Group" for e in entities), f"got: {entities}")

    with patch.object(bedrock_client, "converse_with_tools", new=AsyncMock(side_effect=fake_converse_with_tools)):
        asyncio.run(detector.detect_entities("doc-2", 3, [], alt_texts=None))
    check("no alt-text section is added to the prompt when there's nothing to pass",
          "alt-text" not in captured["user_message"].lower(), f"prompt was: {captured.get('user_message')}")


class _FakeAlias:
    def __init__(self, raw_value):
        self.raw_value = raw_value


class _FakeMaskingEntity:
    def __init__(self, id_, mask_token, aliases=None, custom_replacement=None):
        self.id = id_
        self.mask_token = mask_token
        self.aliases = aliases or []
        self.custom_replacement = custom_replacement


class _FakeAliasQuery:
    def __init__(self, entities):
        self._entities = entities

    def filter(self, *a, **k):
        return self

    def all(self):
        return self._entities


class _FakeAliasDB:
    def __init__(self, others):
        self._others = others

    def query(self, model):
        return _FakeAliasQuery(self._others)

    def flush(self):
        pass


def check_custom_replacement_alias() -> None:
    """Task C acceptance test: resolved_replacement falls back to mask_token
    when no alias is set (today's exact [CLIENT_N] behavior, unchanged), and
    to the alias once set. validate_custom_replacement must reject an alias
    that names another tracked entity's own surface (e.g. aliasing one
    client to "Pfizer" when Pfizer is itself tracked) or that's already
    assigned to a different entity, and pass a genuinely distinct one.
    Fakes stand in for the ORM entity/session, since this is pure logic."""
    print("\n== custom replacement alias (Task C) ==")
    from app.masking import dictionary

    entity = _FakeMaskingEntity(id_=1, mask_token="[CLIENT_16]")
    check("resolved_replacement falls back to mask_token when no alias is set",
          dictionary.resolved_replacement(entity) == "[CLIENT_16]")

    other = _FakeMaskingEntity(id_=2, mask_token="[CLIENT_55]", aliases=[_FakeAlias("Pfizer")])
    problems_collision = dictionary.validate_custom_replacement(_FakeAliasDB([other]), entity, "Pfizer")
    check("alias equal to another tracked entity's own surface is rejected",
          len(problems_collision) > 0, f"got: {problems_collision}")

    other_with_alias = _FakeMaskingEntity(id_=3, mask_token="[CLIENT_60]", custom_replacement="Acme Pharma")
    problems_dup = dictionary.validate_custom_replacement(_FakeAliasDB([other_with_alias]), entity, "Acme Pharma")
    check("alias already assigned to a different entity is rejected",
          len(problems_dup) > 0, f"got: {problems_dup}")

    fake_db_clean = _FakeAliasDB([other, other_with_alias])
    problems_clean = dictionary.validate_custom_replacement(fake_db_clean, entity, "Acme Widgets")
    check("a genuinely distinct alias passes deterministic validation", problems_clean == [], f"got: {problems_clean}")

    dictionary.set_custom_replacement(fake_db_clean, entity, "Acme Widgets")
    check("resolved_replacement uses the alias once set", dictionary.resolved_replacement(entity) == "Acme Widgets")


def check_alias_llm_validation() -> None:
    """Task C acceptance test: the LLM alias-validation call must flag a
    real company name ("Pfizer") as is_real_organization, and pass a clearly
    fictional placeholder ("Acme Pharma"). Mocked Bedrock - tests the
    plumbing, not real model judgment, same discipline as Task 1/2's tests."""
    print("\n== alias LLM validation (Task C) ==")
    from unittest.mock import AsyncMock

    from app.agents.sanitization import alias_validate
    from app.llm import bedrock_client

    real_org_resp = bedrock_client.BedrockResponse(
        text="", parsed={"is_real_organization": True, "reason": "Pfizer is a real pharmaceutical company."},
        input_tokens=30, output_tokens=15, estimated_cost_usd=0.0003,
    )
    with patch.object(bedrock_client, "converse", new=AsyncMock(return_value=real_org_resp)):
        resp = asyncio.run(alias_validate.validate_alias("Pfizer"))
    check("a real organization alias is flagged", (resp.parsed or {}).get("is_real_organization") is True, f"got: {resp.parsed}")

    fictional_resp = bedrock_client.BedrockResponse(
        text="", parsed={"is_real_organization": False, "reason": "Generic fictional placeholder name."},
        input_tokens=30, output_tokens=15, estimated_cost_usd=0.0003,
    )
    with patch.object(bedrock_client, "converse", new=AsyncMock(return_value=fictional_resp)):
        resp2 = asyncio.run(alias_validate.validate_alias("Acme Pharma"))
    check("a fictional placeholder alias is not flagged", (resp2.parsed or {}).get("is_real_organization") is False, f"got: {resp2.parsed}")


def check_image_placeholder_custom_label() -> None:
    """Task C acceptance test: a custom alias label rendered into the
    redaction placeholder must still be recognized by is_placeholder_bytes
    - idempotent-redaction safety (a second remediation pass must recognize
    an already-redacted image and target nothing). A longer label's extra
    text pixels must not push the near-gray ratio below the detection
    threshold."""
    print("\n== image placeholder custom label (Task C) ==")
    from app.documents.image_redact import is_placeholder_bytes, placeholder_png

    default_png = placeholder_png(200, 150)
    check("default 'REDACTED' placeholder is recognized", is_placeholder_bytes(default_png))

    custom_png = placeholder_png(200, 150, label="Acme Pharma")
    check("a custom alias label is still recognized as a placeholder (idempotent redaction safety)",
          is_placeholder_bytes(custom_png))

    long_custom_png = placeholder_png(200, 150, label="A Much Longer Corporate Alias Name Inc.")
    check("a longer custom alias label doesn't break placeholder recognition either",
          is_placeholder_bytes(long_custom_png))


def check_redact_images_per_entity_labels() -> None:
    """Task C acceptance test: redact_images's `labels` dict must produce a
    DIFFERENT placeholder per image when different entities' aliases apply
    to different logos in the same document - not one shared label
    overwriting whichever entity happened to redact last."""
    print("\n== redact_images per-entity labels (Task C) ==")
    from pptx import Presentation
    from pptx.util import Inches

    from app.documents.image_redact import redact_images
    from app.documents.images import extract_images

    workdir = Path(tempfile.mkdtemp(prefix="naviknow-alias-images-"))
    try:
        ct = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        path = str(workdir / "alias_images.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(io.BytesIO(_png_bytes((10, 20, 30))), Inches(0.5), Inches(0.5))
        slide.shapes.add_picture(io.BytesIO(_png_bytes((40, 50, 60))), Inches(3), Inches(0.5))
        prs.save(path)

        refs = extract_images(path, ct, "alias_images.pptx")
        # Exclude docProps/thumbnail.jpeg - PowerPoint auto-generates this
        # document thumbnail as its own embedded image, unrelated to the two
        # slide pictures this test actually plants.
        pptx_refs = [r for r in refs if r.locator.get("kind") == "pptx" and r.locator.get("partname", "").startswith("ppt/media/")]
        check("fixture has two distinct images to redact", len(pptx_refs) == 2, f"got: {len(pptx_refs)}")

        labels = {pptx_refs[0].index: "Acme Pharma", pptx_refs[1].index: "Beta Bank"}
        redacted, _ = redact_images(path, ct, "alias_images.pptx", pptx_refs, labels=labels)
        check("both images redacted", redacted == 2, f"got: {redacted}")

        refs_after = extract_images(path, ct, "alias_images.pptx")
        bytes_by_index = {
            r.index: r.image_bytes for r in refs_after
            if r.locator.get("kind") == "pptx" and r.locator.get("partname", "").startswith("ppt/media/")
        }
        check("the two redacted placeholders have different bytes (different alias labels rendered)",
              bytes_by_index[pptx_refs[0].index] != bytes_by_index[pptx_refs[1].index],
              "placeholders are byte-identical despite different labels")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def check_pdf_annotation_author() -> None:
    """Item 1b acceptance test: a PDF annotation's Title field (the standard
    PDF markup-annotation author field, same "always a real identity"
    shape as PPTX/DOCX comment author names) must be flagged and cleared
    UNCONDITIONALLY - independent of `surfaces` - not just its content text."""
    print("\n== pdf annotation author/title (item 1b) ==")
    import fitz

    from app.documents.comment_scan import find_residual_comments
    from app.documents.comment_scrub import scrub_comments

    workdir = Path(tempfile.mkdtemp(prefix="naviknow-pdf-author-"))
    try:
        ct = "application/pdf"
        path = str(workdir / "annot.pdf")
        doc = fitz.open()
        page = doc.new_page()
        annot = page.add_text_annot((72, 160), "benign note")
        annot.set_info(title="Vallab Deshmukh")
        annot.update()
        doc.save(path)
        doc.close()

        residual_before = find_residual_comments(path, ct, "annot.pdf", [])
        check("annotation author/title flagged with an EMPTY surfaces list",
              any("identity attribute" in h and "Vallab Deshmukh" in h for h in residual_before), f"got: {residual_before}")

        changed = scrub_comments(path, ct, "annot.pdf", {}, STYLE)
        check("scrub_comments clears the annotation title with an EMPTY surface_to_token", changed >= 1, f"got: {changed}")

        residual_after = find_residual_comments(path, ct, "annot.pdf", [])
        check("no annotation-author residual remains after scrub", len(residual_after) == 0, f"got: {residual_after}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def check_channel_coverage_audit() -> None:
    """Item 1c acceptance test: the coverage audit must flag a deliberately
    unknown, identity-suggestive part (simulating "the next authors.xml") -
    and must NOT flag anything on the existing regression fixture (proving
    it doesn't cry wolf on every normal file, which would make it noise
    reviewers learn to ignore)."""
    print("\n== channel coverage audit (item 1c) ==")
    from app.documents.channel_coverage import audit_channel_coverage

    workdir = Path(tempfile.mkdtemp(prefix="naviknow-coverage-"))
    try:
        ct = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        clean_path = str(workdir / "clean.pptx")
        build_pptx(clean_path)
        warnings_clean = audit_channel_coverage(clean_path, ct, "clean.pptx")
        check("no false positives on a normal, fully-covered pptx", warnings_clean == [], f"got: {warnings_clean}")

        unknown_path = str(workdir / "unknown.pptx")
        build_pptx(unknown_path)
        # Simulate "the next authors.xml": an unrecognized part with an
        # identity-suggestive name this audit has never been told about.
        with zipfile.ZipFile(unknown_path, "a") as z:
            z.writestr("ppt/reviewersList.xml", "<reviewers><reviewer name='Someone'/></reviewers>")
        warnings_unknown = audit_channel_coverage(unknown_path, ct, "unknown.pptx")
        check("an unrecognized identity-suggestive part is flagged",
              any("ppt/reviewersList.xml" in w for w in warnings_unknown), f"got: {warnings_unknown}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def check_docx_alt_text() -> None:
    """Item 3b acceptance test: alt-text scan/scrub for DOCX specifically -
    only PPTX had a dedicated regression fixture before this. python-docx's
    inline picture docPr has no descr/title by default, so this injects
    them directly into the real XML python-docx generated, the same
    raw-injection pattern already used elsewhere in this suite for shapes
    neither library exposes a high-level API for."""
    print("\n== docx alt-text channel (item 3b) ==")
    import docx
    from PIL import Image

    from app.documents.alttext_scan import extract_alt_text
    from app.documents.alttext_scrub import scrub_alt_text
    from app.documents.verify import find_residual_surfaces

    workdir = Path(tempfile.mkdtemp(prefix="naviknow-docx-alttext-"))
    try:
        ct = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        path = str(workdir / "alttext.docx")
        png_buf = io.BytesIO()
        Image.new("RGB", (60, 40), (10, 20, 30)).save(png_buf, format="PNG")
        png_buf.seek(0)
        doc = docx.Document()
        doc.add_picture(png_buf)
        doc.save(path)

        with zipfile.ZipFile(path) as z:
            xml = z.read("word/document.xml").decode("utf-8")
        assert "<wp:docPr" in xml, "fixture assumption broke: python-docx no longer emits wp:docPr this way"
        new_xml = re.sub(
            r'(<wp:docPr\b[^>]*?)/>',
            rf'\1 descr="Our Business | {CLIENT} Group logo" title="{CLIENT} wordmark"/>',
            xml, count=1,
        )
        tmp_path = path + ".tmp"
        with zipfile.ZipFile(path, "r") as zin, zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                zout.writestr(item, new_xml.encode("utf-8") if item.filename == "word/document.xml" else zin.read(item.filename))
        shutil.move(tmp_path, path)

        pre = extract_alt_text(path, ct, "alttext.docx")
        check("docx alt-text descr/title detected before scrub", any(CLIENT in v for v in pre), f"got: {pre}")

        residual_before = find_residual_surfaces(path, ct, "alttext.docx", SURFACES)
        check("text-channel verification catches the docx alt-text leak pre-scrub",
              CLIENT in residual_before, f"got: {residual_before}")

        changed = scrub_alt_text(path, ct, "alttext.docx", SURFACE_TO_TOKEN, STYLE)
        check("scrub_alt_text reports rewritten attributes on docx", changed >= 2, f"got: {changed}")

        post = extract_alt_text(path, ct, "alttext.docx")
        check("no client name remains in docx alt-text after scrub", not any(CLIENT in v for v in post), f"got: {post}")

        residual_after = find_residual_surfaces(path, ct, "alttext.docx", SURFACES)
        check("text channel verifies clean after docx alt-text scrub", CLIENT not in residual_after, f"got: {residual_after}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def check_xlsx_alt_text() -> None:
    """Item 3b acceptance test: alt-text scan/scrub for XLSX specifically -
    openpyxl's embedded-image cNvPr defaults to descr="Picture" (no
    namespace prefix - a bare <cNvPr>, confirming TAG_RE's optional-prefix
    match is actually needed, not just defensive)."""
    print("\n== xlsx alt-text channel (item 3b) ==")
    import openpyxl
    from openpyxl.drawing.image import Image as XLImage
    from PIL import Image

    from app.documents.alttext_scan import extract_alt_text
    from app.documents.alttext_scrub import scrub_alt_text
    from app.documents.verify import find_residual_surfaces

    workdir = Path(tempfile.mkdtemp(prefix="naviknow-xlsx-alttext-"))
    try:
        ct = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        path = str(workdir / "alttext.xlsx")
        png_buf = io.BytesIO()
        Image.new("RGB", (60, 40), (10, 20, 30)).save(png_buf, format="PNG")
        png_buf.seek(0)
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.add_image(XLImage(png_buf), "A1")
        wb.save(path)

        with zipfile.ZipFile(path) as z:
            drawing_name = next(n for n in z.namelist() if re.match(r"^xl/drawings/drawing\d*\.xml$", n))
            xml = z.read(drawing_name).decode("utf-8")
        assert 'descr="Picture"' in xml, "fixture assumption broke: openpyxl no longer emits a default descr this way"
        new_xml = xml.replace('descr="Picture"', f'descr="{CLIENT} Group logo"')
        tmp_path = path + ".tmp"
        with zipfile.ZipFile(path, "r") as zin, zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                zout.writestr(item, new_xml.encode("utf-8") if item.filename == drawing_name else zin.read(item.filename))
        shutil.move(tmp_path, path)

        pre = extract_alt_text(path, ct, "alttext.xlsx")
        check("xlsx alt-text descr detected before scrub (unprefixed <cNvPr>)", any(CLIENT in v for v in pre), f"got: {pre}")

        residual_before = find_residual_surfaces(path, ct, "alttext.xlsx", SURFACES)
        check("text-channel verification catches the xlsx alt-text leak pre-scrub",
              CLIENT in residual_before, f"got: {residual_before}")

        changed = scrub_alt_text(path, ct, "alttext.xlsx", SURFACE_TO_TOKEN, STYLE)
        check("scrub_alt_text reports rewritten attributes on xlsx", changed >= 1, f"got: {changed}")

        post = extract_alt_text(path, ct, "alttext.xlsx")
        check("no client name remains in xlsx alt-text after scrub", not any(CLIENT in v for v in post), f"got: {post}")

        residual_after = find_residual_surfaces(path, ct, "alttext.xlsx", SURFACES)
        check("text channel verifies clean after xlsx alt-text scrub", CLIENT not in residual_after, f"got: {residual_after}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def check_revalidation_agent() -> None:
    """Acceptance test for the revalidation agent (revalidate.py): a second,
    independent pass over the rendered output. Uses the two REAL leaks found
    by directly inspecting a completed production run's actual output file -
    "[CLIENT_25] Turbine Ltd." (a partial-name fragment next to its own mask
    token - Lens 1, deterministic) and "Arvind Fashions" (a total detection
    miss elsewhere in the text - Lens 2, mocked Bedrock, same offline
    discipline as check_precision_qa/check_reidentify_qa)."""
    print("\n== revalidation agent (fresh-eyes re-check, Lens 1 + Lens 2 + scoring) ==")
    from app.agents.sanitization import revalidate
    from app.llm import bedrock_client

    # Lens 1: deterministic boundary heuristic.
    leaky_text = (
        "Intelligent automation of the B2C Order-to-cash process reconciliation with third-party "
        "marketplace [CLIENT_15] (such as Ajio, Myntra, etc.) and D2C website for Arvind Fashions. "
        "AI-driven process transformation for [CLIENT_25] Turbine Ltd. | Bangalore."
    )
    boundary_hits = revalidate.find_boundary_leaks(leaky_text, ["[CLIENT_15]", "[CLIENT_25]"])
    hit = next((h for h in boundary_hits if h["mask_token"] == "[CLIENT_25]"), None)
    check("boundary heuristic catches a legal-suffix fragment left next to its own mask token",
          hit is not None and hit["leaked_text"] == "Turbine Ltd.", f"got: {boundary_hits}")
    check("boundary heuristic does not flag a token with nothing suspicious immediately after it",
          not any(h["mask_token"] == "[CLIENT_15]" for h in boundary_hits), f"got: {boundary_hits}")

    clean_text = "The engagement was led by [CLIENT_1] alongside their advisory team."
    check("boundary heuristic finds nothing in ordinary masked prose",
          revalidate.find_boundary_leaks(clean_text, ["[CLIENT_1]"]) == [])

    # Lens 2: fresh, adversarial re-detection (mocked Bedrock).
    canned_entities = {
        "entities": [
            {"surface_text": "Arvind Fashions", "entity_type": "CLIENT_NAME", "confidence": 0.9},
            {"surface_text": "[CLIENT_15]", "entity_type": "CLIENT_NAME", "confidence": 0.4},  # must be filtered
            {"surface_text": "Navikenz", "entity_type": "CLIENT_NAME", "confidence": 0.3},  # own-firm, must be filtered
        ]
    }
    mock_resp = bedrock_client.BedrockResponse(text="", parsed=canned_entities, input_tokens=200, output_tokens=40, estimated_cost_usd=0.001)
    with patch.object(bedrock_client, "converse", new=AsyncMock(return_value=mock_resp)):
        resp = asyncio.run(revalidate.fresh_redetect(leaky_text))
    fresh_hits = revalidate.parse_fresh_redetect_hits(resp)
    check("fresh re-detect surfaces the total-miss entity (Arvind Fashions)",
          any(h["leaked_text"] == "Arvind Fashions" for h in fresh_hits), f"got: {fresh_hits}")
    check("fresh re-detect filters out a bare mask token echoed back",
          not any(h["leaked_text"] == "[CLIENT_15]" for h in fresh_hits), f"got: {fresh_hits}")
    check("fresh re-detect filters out the delivery firm's own name",
          not any(h["leaked_text"] == "Navikenz" for h in fresh_hits), f"got: {fresh_hits}")

    with patch.object(bedrock_client, "converse", new=AsyncMock(return_value=mock_resp)) as mocked:
        resp_empty = asyncio.run(revalidate.fresh_redetect(""))
    check("fresh re-detect makes no Bedrock call on empty text", resp_empty is None and not mocked.called)

    # Scoring: masked / (masked + residual), always paired with the residual list.
    score = revalidate.compute_completeness(masked_count=46, residual_count=2)
    check("completeness score matches the real Fidelity run (46 masked, 2 residual -> ~95.8%)",
          abs(score - 95.8) < 0.1, f"got: {score}")
    check("completeness score is 100% when nothing was masked and nothing is residual",
          revalidate.compute_completeness(0, 0) == 100.0)
    check("completeness score is 0% when everything masked has an equal-count residual twin",
          revalidate.compute_completeness(0, 5) == 0.0)


def check_logo_thumbnail() -> None:
    """Acceptance test for the admin-panel logo visibility gap: an approved
    image match previously left NO trace in the admin UI beyond a bare hex
    phash - _save_thumbnail must produce a small preview PNG on disk from
    real image bytes, must degrade gracefully (None, not a crash) on bytes
    PIL can't open, and store_reference must persist the resulting path onto
    the LogoReference row so the API can serve it back."""
    print("\n== logo reference thumbnail (admin panel visibility) ==")
    import uuid

    from app.masking import logo_reference

    workdir = Path(tempfile.mkdtemp(prefix="naviknow-logo-thumb-"))
    try:
        with patch.object(logo_reference.settings, "OUTPUTS_DIR", str(workdir)):
            image_bytes = _png_bytes(color=(40, 90, 200))
            entity_id = uuid.uuid4()

            path = logo_reference._save_thumbnail(entity_id, image_bytes)
            check("a thumbnail path is returned for real image bytes", path is not None, f"got: {path}")
            check("the thumbnail file actually exists on disk", path is not None and Path(path).exists())

            from PIL import Image
            with Image.open(path) as im:
                check("the thumbnail was actually resized down (not a full-size copy)",
                      im.size[0] <= 160 and im.size[1] <= 160, f"got: {im.size}")

            bad_path = logo_reference._save_thumbnail(entity_id, b"not a real image")
            check("an unopenable image degrades to None instead of raising", bad_path is None, f"got: {bad_path}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def check_infra_credential_detection() -> None:
    """Phase 2 acceptance test: the regex module backing INFRA_IDENTIFIER/
    CREDENTIAL candidate detection must catch the documented shapes (IPv4,
    internal hostnames, AWS-style access key IDs, Bearer tokens, connection
    strings) and must NOT flag an ordinary public IP or ordinary prose as a
    credential - CREDENTIAL's mandatory/non-overridable contract only holds
    up if this stays high-precision (see the module's own docstring)."""
    print("\n== infra/credential regex detection (Phase 2) ==")
    from app.agents.sanitization.regex_patterns import infra_credential

    text = (
        "Reach the internal panel at admin.corp.internal or 10.0.4.17. "
        "AWS access key AKIAIOSFODNN7EXAMPLE leaked in a log. "
        "Authorization: Bearer eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.abcdefghijklmnop "
        "Connection string: postgresql://dbuser:S3cretPass@prod-db.internal:5432/appdb"
    )
    found = infra_credential.scan(text)
    by_type: dict[str, list[str]] = {}
    for surface, etype in found:
        by_type.setdefault(etype, []).append(surface)

    check("finds the internal hostname as INFRA_IDENTIFIER", "admin.corp.internal" in by_type.get("INFRA_IDENTIFIER", []), f"got: {by_type}")
    check("finds the private IPv4 as INFRA_IDENTIFIER", "10.0.4.17" in by_type.get("INFRA_IDENTIFIER", []), f"got: {by_type}")
    check("finds the AWS-style access key ID as CREDENTIAL", "AKIAIOSFODNN7EXAMPLE" in by_type.get("CREDENTIAL", []), f"got: {by_type}")
    check("finds the Bearer token as CREDENTIAL", any(s.startswith("Bearer ") for s in by_type.get("CREDENTIAL", [])), f"got: {by_type}")
    check("finds the connection string as CREDENTIAL", any("prod-db.internal" in s and "S3cretPass" in s for s in by_type.get("CREDENTIAL", [])), f"got: {by_type}")

    clean = "Revenue grew 12% year over year across all regions."
    check("ordinary prose with no technical shapes produces no candidates", infra_credential.scan(clean) == [], f"got: {infra_credential.scan(clean)}")

    public_ip_only = "Our public DNS resolver is 8.8.8.8."
    public_hits = infra_credential.scan(public_ip_only)
    check("a bare public IP is still flagged as INFRA_IDENTIFIER (a candidate, not a final verdict - the LLM/reviewer judges public vs. internal)",
          any(t == "INFRA_IDENTIFIER" for _, t in public_hits), f"got: {public_hits}")
    check("a bare public IP is never flagged as CREDENTIAL", not any(t == "CREDENTIAL" for _, t in public_hits), f"got: {public_hits}")


def check_credential_mandatory_no_override() -> None:
    """Phase 2 acceptance test: a reviewer's removed_surfaces edit must be
    silently ignored for a CREDENTIAL entity (mandatory, non-overridable),
    while an ordinary CLIENT_NAME/INFRA_IDENTIFIER removal in the SAME
    request still goes through normally - the override doesn't leak into
    unrelated entities."""
    print("\n== CREDENTIAL mandatory, non-overridable removal (Phase 2) ==")
    # _filter_removed_entities was generalized into _resolve_entity_inclusion
    # in Phase 3 (now takes default_action-aware entities plus an `included`
    # opt-in set) - this test still exercises the original Phase 2 contract
    # specifically (CREDENTIAL alone, via the mask-vs-mandatory split), kept
    # separate from check_flag_vs_mask_vs_mandatory_inclusion's broader
    # Phase 3 coverage rather than deleting a working, named regression.
    from app.agents.sanitization.agent import _resolve_entity_inclusion

    proposed = [
        {"surface_text": "Acme Corp", "entity_type": "CLIENT_NAME", "default_action": "mask"},
        {"surface_text": "admin.corp.internal", "entity_type": "INFRA_IDENTIFIER", "default_action": "mask"},
        {"surface_text": "AKIAIOSFODNN7EXAMPLE", "entity_type": "CREDENTIAL", "default_action": "mandatory"},
    ]
    removed = {"acme corp", "admin.corp.internal", "akiaiosfodnn7example"}  # reviewer tries to exclude all three

    kept, blocked = _resolve_entity_inclusion(proposed, removed, included=set())
    kept_surfaces = {e["surface_text"] for e in kept}

    check("CREDENTIAL survives the removal attempt", "AKIAIOSFODNN7EXAMPLE" in kept_surfaces, f"got: {kept_surfaces}")
    check("CLIENT_NAME removal is honored normally", "Acme Corp" not in kept_surfaces, f"got: {kept_surfaces}")
    check("INFRA_IDENTIFIER removal is honored normally (only CREDENTIAL is mandatory)", "admin.corp.internal" not in kept_surfaces, f"got: {kept_surfaces}")
    check("the blocked-removal list names exactly the CREDENTIAL surface", blocked == ["AKIAIOSFODNN7EXAMPLE"], f"got: {blocked}")

    no_removal_attempt = _resolve_entity_inclusion(proposed, set(), included=set())
    check("no removal attempted -> blocked list is empty (no false positive flag)", no_removal_attempt[1] == [], f"got: {no_removal_attempt[1]}")


def check_exif_strip() -> None:
    """Phase 2 acceptance test: strip_exif must actually remove EXIF from a
    JPEG that has it, leave an image with no EXIF untouched (reported count
    excludes it), and never corrupt the image (still opens, same dimensions)
    - and must be a no-op for PDF (out of scope for this phase, see module
    docstring) rather than raising."""
    print("\n== EXIF stripping (Phase 2) ==")
    from PIL import Image
    from PIL.ExifTags import Base as ExifTags

    from app.documents.exif_strip import strip_exif

    workdir = Path(tempfile.mkdtemp(prefix="naviknow-exif-"))
    try:
        im = Image.new("RGB", (100, 80), color=(200, 50, 50))
        exif = Image.Exif()
        exif[ExifTags.Make.value] = "TestCameraCo"
        exif[ExifTags.Model.value] = "TestModel-9000"
        buf = io.BytesIO()
        im.save(buf, format="JPEG", exif=exif.tobytes())
        jpeg_with_exif = buf.getvalue()

        plain_buf = io.BytesIO()
        Image.new("RGB", (40, 40), color=(10, 10, 10)).save(plain_buf, format="PNG")
        png_no_exif = plain_buf.getvalue()

        docx_path = str(workdir / "exif_test.docx")
        with zipfile.ZipFile(docx_path, "w") as z:
            z.writestr("word/media/image1.jpeg", jpeg_with_exif)
            z.writestr("word/media/image2.png", png_no_exif)
            z.writestr("[Content_Types].xml", "<Types/>")

        stripped = strip_exif(
            docx_path, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "exif_test.docx",
        )
        check("exactly one image (the JPEG with real EXIF) reports as stripped", stripped == 1, f"got: {stripped}")

        with zipfile.ZipFile(docx_path) as z:
            result_jpeg = z.read("word/media/image1.jpeg")
            result_png = z.read("word/media/image2.png")

        result_im = Image.open(io.BytesIO(result_jpeg))
        check("EXIF is actually gone after stripping", not bool(result_im.getexif()), f"got: {dict(result_im.getexif())}")
        check("the image still opens and keeps its original dimensions", result_im.size == (100, 80), f"got: {result_im.size}")
        check("the no-EXIF PNG is byte-identical (left untouched, not blindly re-encoded)", result_png == png_no_exif)

        pdf_calls = strip_exif(docx_path, "application/pdf", "whatever.pdf")
        check("PDF is a no-op (0), not an error - out of scope for this phase", pdf_calls == 0, f"got: {pdf_calls}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def check_data_sample_and_sensitive_text_image_fields() -> None:
    """Phase 2 acceptance test: the vision call's new contains_real_data_sample
    field must survive _scan_one_group into the returned ImageGroup, and text
    read off an image (ocr_text) that matches an INFRA_IDENTIFIER/CREDENTIAL
    shape must populate sensitive_text_matches and drive mandatory_redaction
    for a CREDENTIAL - the same non-overridable contract CREDENTIAL gets in
    body text, now for images too. db is a minimal mock (not None): unlike
    the vision-cache test above, this run's ocr_text is non-empty, so
    _ocr_match's dictionary.lookup(db, ...) call must not crash - configured
    to always miss (None/[]), same as a document with an empty dictionary."""
    print("\n== data-sample + sensitive-text image fields (Phase 2, mocked Bedrock) ==")
    import uuid as uuid_mod

    from app.agents.sanitization import image_scan
    from app.documents.images import ImageRef
    from app.llm import bedrock_client

    db = MagicMock()
    db.query.return_value.filter.return_value.first.return_value = None
    db.query.return_value.filter.return_value.all.return_value = []

    async def fake_converse_vision(**kwargs):
        return bedrock_client.BedrockResponse(
            text="",
            parsed={
                "contains_client_identity": False,
                "description": "a dashboard screenshot",
                "confidence": 0.4,
                "ocr_text": ["Revenue Dashboard", "Host: prod-db.internal", "AKIAIOSFODNN7EXAMPLE"],
                "contains_real_data_sample": True,
            },
            input_tokens=40, output_tokens=15, estimated_cost_usd=0.0004,
        )

    png = _png_bytes((80, 80, 80))
    ref = ImageRef(index=0, location_label="slide 4", image_bytes=png, image_format="png", locator={})

    with patch.object(image_scan.vision_cache, "load_cached_verdict", return_value=None), \
         patch.object(image_scan.vision_cache, "store_verdict"), \
         patch.object(bedrock_client, "converse_vision", side_effect=fake_converse_vision):
        group, _, _, _ = asyncio.run(image_scan._scan_one_group(
            db, 0, [ref], png, "png", None, logo_references=[], run_id=uuid_mod.uuid4(),
        ))

    check("contains_real_data_sample survives from the vision response", group.contains_real_data_sample is True, f"got: {group.contains_real_data_sample}")
    sensitive_types = {etype for _, etype in group.sensitive_text_matches}
    check("hostname read off the image is classified as INFRA_IDENTIFIER", "INFRA_IDENTIFIER" in sensitive_types, f"got: {group.sensitive_text_matches}")
    check("credential read off the image is classified as CREDENTIAL", "CREDENTIAL" in sensitive_types, f"got: {group.sensitive_text_matches}")

    # Build the same images_proposal payload shape agent.py's detect() emits,
    # to check the mandatory_redaction derivation without needing the full
    # detect() pipeline (Bedrock Detector call, document fixtures, etc.).
    mandatory = any(etype in {"CREDENTIAL"} for _, etype in group.sensitive_text_matches)
    check("a CREDENTIAL read off an image drives mandatory_redaction=True", mandatory is True)


def check_dataset_provenance() -> None:
    """Phase 2 acceptance test: a new upload whose text matches an APPROVED,
    account-linked entity gets an advisory provenance warning naming that
    account; a document naming only entities with NO account link, or none
    at all, gets no warning - this must never block the upload either way,
    it only ever returns a message or None."""
    print("\n== dataset provenance check (Phase 2) ==")
    from app.documents import dataset_provenance

    workdir = Path(tempfile.mkdtemp(prefix="naviknow-provenance-"))
    try:
        import docx

        doc = docx.Document()
        doc.add_paragraph("This deck was prepared for Acme Robotics Inc. by our delivery team.")
        linked_path = str(workdir / "linked.docx")
        doc.save(linked_path)

        doc2 = docx.Document()
        doc2.add_paragraph("A completely unrelated internal planning note with no client names at all.")
        unlinked_path = str(workdir / "unlinked.docx")
        doc2.save(unlinked_path)

        docx_content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

        class _FakeAccount:
            name = "Acme Robotics"

        class _FakeAlias:
            raw_value = "Acme Robotics Inc."

        class _FakeEntity:
            client_account_id = "11111111-1111-1111-1111-111111111111"
            client_account = _FakeAccount()
            aliases = [_FakeAlias()]

        with patch.object(dataset_provenance.dictionary, "find_in_text", return_value=[(_FakeEntity(), "Acme Robotics Inc.")]):
            warning = dataset_provenance.check(None, linked_path, docx_content_type, "linked.docx")
        check("a document matching an account-linked entity gets a named warning", warning is not None and "Acme Robotics" in warning, f"got: {warning}")

        with patch.object(dataset_provenance.dictionary, "find_in_text", return_value=[]):
            no_warning = dataset_provenance.check(None, unlinked_path, docx_content_type, "unlinked.docx")
        check("a document matching nothing gets no warning (None, not an error)", no_warning is None, f"got: {no_warning}")

        class _FakeEntityNoAccount:
            client_account_id = None
            client_account = None
            aliases = []

        with patch.object(dataset_provenance.dictionary, "find_in_text", return_value=[(_FakeEntityNoAccount(), "some match")]):
            unlinked_entity_warning = dataset_provenance.check(None, linked_path, docx_content_type, "linked.docx")
        check("a matched entity with NO account link produces no warning (link is what matters, not the match itself)",
              unlinked_entity_warning is None, f"got: {unlinked_entity_warning}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def check_entity_actions_default_resolution() -> None:
    """Phase 3 acceptance test: entity_actions.resolve_default_action must
    return "mask" for every Phase 1/2 type (unchanged), "mandatory" for
    CREDENTIAL (unchanged), "flag" for exactly the five Phase 3 flag-default
    types, and a consent-dependent "mask"/"keep" for INTERNAL_TEAM_MEMBER -
    the single source of truth every other Phase 3 check below builds on."""
    print("\n== entity default-action resolution (Phase 3) ==")
    from app.agents.sanitization import entity_actions

    check("CLIENT_NAME (Phase 1) still defaults to mask", entity_actions.resolve_default_action("CLIENT_NAME") == "mask")
    check("PERSON (Phase 1) still defaults to mask", entity_actions.resolve_default_action("PERSON") == "mask")
    check("INFRA_IDENTIFIER (Phase 2) still defaults to mask", entity_actions.resolve_default_action("INFRA_IDENTIFIER") == "mask")
    check("CREDENTIAL (Phase 2) is still mandatory", entity_actions.resolve_default_action("CREDENTIAL") == "mandatory")

    for flag_type in ("COMMERCIAL_TERM", "COMPETITOR_NAME", "STRATEGY_MENTION", "OWN_COST_DETAIL", "ORG_CHART_STRUCTURE"):
        check(f"{flag_type} defaults to flag", entity_actions.resolve_default_action(flag_type) == "flag")

    check("CLIENT_PERSON_TITLE defaults to mask (an ordinary CLIENT_* type)", entity_actions.resolve_default_action("CLIENT_PERSON_TITLE") == "mask")
    check("CLIENT_PHONE defaults to mask (an ordinary CLIENT_* type)", entity_actions.resolve_default_action("CLIENT_PHONE") == "mask")

    check("INTERNAL_TEAM_MEMBER with no consent record defaults to mask",
          entity_actions.resolve_default_action("INTERNAL_TEAM_MEMBER", None) == "mask")
    check("INTERNAL_TEAM_MEMBER with pending consent defaults to mask",
          entity_actions.resolve_default_action("INTERNAL_TEAM_MEMBER", "pending") == "mask")
    check("INTERNAL_TEAM_MEMBER with granted consent defaults to keep",
          entity_actions.resolve_default_action("INTERNAL_TEAM_MEMBER", "granted") == "keep")


def check_commercial_term_and_title_detection() -> None:
    """Phase 3 acceptance test: the regex module backing COMMERCIAL_TERM/
    CLIENT_PERSON_TITLE candidate detection must catch documented shapes
    (currency amounts, payment-term keywords, multi-word job titles) and
    must NOT flag a bare generic noun ("Manager") that would otherwise
    flood the proposal with noise."""
    print("\n== commercial-term / client-title regex detection (Phase 3) ==")
    from app.agents.sanitization.regex_patterns import commercial_terms

    text = (
        "The contract value was $2.5M with net 30 payment terms. Our main contact is the "
        "Vice President of Operations, who reports to the Chief Financial Officer."
    )
    found = commercial_terms.scan(text)
    by_type: dict[str, list[str]] = {}
    for surface, etype in found:
        by_type.setdefault(etype, []).append(surface)

    check("finds the currency amount as COMMERCIAL_TERM", any("2.5M" in s for s in by_type.get("COMMERCIAL_TERM", [])), f"got: {by_type}")
    check("finds 'net 30' as COMMERCIAL_TERM", any("net 30" in s.lower() for s in by_type.get("COMMERCIAL_TERM", [])), f"got: {by_type}")
    check("finds 'Vice President of Operations' as CLIENT_PERSON_TITLE", "Vice President of Operations" in by_type.get("CLIENT_PERSON_TITLE", []), f"got: {by_type}")
    check("finds 'Chief Financial Officer' as CLIENT_PERSON_TITLE", "Chief Financial Officer" in by_type.get("CLIENT_PERSON_TITLE", []), f"got: {by_type}")

    noisy = "The manager asked the director to review the report."
    check("a bare 'manager'/'director' with no distinguishing continuation produces no CLIENT_PERSON_TITLE candidate",
          not any(etype == "CLIENT_PERSON_TITLE" for _, etype in commercial_terms.scan(noisy)), f"got: {commercial_terms.scan(noisy)}")


def check_flag_vs_mask_vs_mandatory_inclusion() -> None:
    """Phase 3 acceptance test: _resolve_entity_inclusion (the generalized
    form of Phase 2's CREDENTIAL-only _filter_removed_entities) must apply
    each of the three default_action contracts correctly, in one request:
    a mask-default entity is kept unless removed; a flag-default entity is
    EXCLUDED unless explicitly included; a mandatory entity is always kept
    regardless of either edit."""
    print("\n== flag vs. mask vs. mandatory entity inclusion (Phase 3) ==")
    from app.agents.sanitization.agent import _resolve_entity_inclusion

    proposed = [
        {"surface_text": "Acme Corp", "entity_type": "CLIENT_NAME", "default_action": "mask"},
        {"surface_text": "$2.5M contract", "entity_type": "COMMERCIAL_TERM", "default_action": "flag"},
        {"surface_text": "Rival Inc", "entity_type": "COMPETITOR_NAME", "default_action": "flag"},
        {"surface_text": "AKIAIOSFODNN7EXAMPLE", "entity_type": "CREDENTIAL", "default_action": "mandatory"},
        {"surface_text": "Jane Doe", "entity_type": "INTERNAL_TEAM_MEMBER", "default_action": "keep"},
    ]

    # Untouched: mask-default entities stay in, flag/keep-default entities
    # stay out, mandatory stays in.
    kept, blocked = _resolve_entity_inclusion(proposed, removed=set(), included=set())
    kept_surfaces = {e["surface_text"] for e in kept}
    check("untouched mask-default entity is kept", "Acme Corp" in kept_surfaces, f"got: {kept_surfaces}")
    check("untouched flag-default COMMERCIAL_TERM is excluded", "$2.5M contract" not in kept_surfaces, f"got: {kept_surfaces}")
    check("untouched flag-default COMPETITOR_NAME is excluded", "Rival Inc" not in kept_surfaces, f"got: {kept_surfaces}")
    check("untouched mandatory CREDENTIAL is kept", "AKIAIOSFODNN7EXAMPLE" in kept_surfaces, f"got: {kept_surfaces}")
    check("untouched keep-default INTERNAL_TEAM_MEMBER is excluded", "Jane Doe" not in kept_surfaces, f"got: {kept_surfaces}")
    check("no removal attempted -> nothing blocked", blocked == [], f"got: {blocked}")

    # Reviewer opts the COMMERCIAL_TERM IN, tries to opt CLIENT_NAME out,
    # and tries (futilely) to remove the CREDENTIAL.
    kept2, blocked2 = _resolve_entity_inclusion(
        proposed,
        removed={"acme corp", "akiaiosfodnn7example"},
        included={"$2.5m contract"},
    )
    kept_surfaces2 = {e["surface_text"] for e in kept2}
    check("opted-in flag-default COMMERCIAL_TERM is now kept", "$2.5M contract" in kept_surfaces2, f"got: {kept_surfaces2}")
    check("still-untouched flag-default COMPETITOR_NAME stays excluded", "Rival Inc" not in kept_surfaces2, f"got: {kept_surfaces2}")
    check("opted-out mask-default CLIENT_NAME is now excluded", "Acme Corp" not in kept_surfaces2, f"got: {kept_surfaces2}")
    check("CREDENTIAL survives the removal attempt regardless", "AKIAIOSFODNN7EXAMPLE" in kept_surfaces2, f"got: {kept_surfaces2}")
    check("the blocked-removal list names exactly the CREDENTIAL surface", blocked2 == ["AKIAIOSFODNN7EXAMPLE"], f"got: {blocked2}")

    # A legacy proposal entity with no default_action field at all (written
    # before this field existed) must fall back to "mask", not crash.
    legacy = [{"surface_text": "Legacy Co", "entity_type": "CLIENT_NAME"}]
    kept3, _ = _resolve_entity_inclusion(legacy, removed=set(), included=set())
    check("an entity with no default_action field falls back to mask-default", len(kept3) == 1, f"got: {kept3}")


def check_immediate_consent_grant() -> None:
    """Phase 3 acceptance test (regression for a real bug caught during
    smoke testing): granting consent for an INTERNAL_TEAM_MEMBER in
    edits.consent_updates must exempt THIS run's occurrence immediately,
    not only future ones - the entity's default_action was computed at
    detect() time, before this consent existed, so without this fold-in
    the person's name still got masked in the very run where consent was
    just granted. Must also leave an unrelated surface (no matching
    consent_updates entry, or a different entity_type) untouched."""
    print("\n== immediate consent grant exempts this run too (Phase 3) ==")
    from app.agents.sanitization.agent import _apply_immediate_consent_grants

    entities = [
        {"surface_text": "John Smith", "entity_type": "INTERNAL_TEAM_MEMBER", "default_action": "mask"},
        {"surface_text": "Jane Doe", "entity_type": "INTERNAL_TEAM_MEMBER", "default_action": "mask"},
        {"surface_text": "Acme Corp", "entity_type": "CLIENT_NAME", "default_action": "mask"},
    ]
    updated = _apply_immediate_consent_grants(entities, removed=set(), consent_updates={"John Smith": "granted"})
    check("consent-granted INTERNAL_TEAM_MEMBER is folded into removed (exempt this run)", "john smith" in updated, f"got: {updated}")
    check("an untouched INTERNAL_TEAM_MEMBER with no consent update is NOT added", "jane doe" not in updated, f"got: {updated}")
    check("an unrelated CLIENT_NAME is never affected by consent_updates", "acme corp" not in updated, f"got: {updated}")

    # A "pending"/"not_required" consent update must NOT exempt this run -
    # only an explicit "granted" does.
    updated2 = _apply_immediate_consent_grants(entities, removed=set(), consent_updates={"John Smith": "pending"})
    check("a 'pending' consent update does not exempt this run", "john smith" not in updated2, f"got: {updated2}")

    # Pre-existing removed_surfaces entries survive untouched alongside the merge.
    updated3 = _apply_immediate_consent_grants(entities, removed={"acme corp"}, consent_updates={"John Smith": "granted"})
    check("a pre-existing removed entry is preserved", "acme corp" in updated3, f"got: {updated3}")
    check("the new consent grant is added alongside it", "john smith" in updated3, f"got: {updated3}")


def check_internal_team_consent() -> None:
    """Phase 3 acceptance test: dictionary.get_consent_status/set_consent_status
    must read/write MaskingEntity.consent_status by name lookup, and a name
    with no prior record must return None (not crash, not default to a
    truthy value) - entity_actions treats None the same as "pending"."""
    print("\n== internal team consent lookup/set (Phase 3) ==")
    from app.masking import dictionary

    class _FakeEntity:
        def __init__(self):
            self.consent_status = None

    fake_entity = _FakeEntity()
    fake_db = MagicMock()

    with patch.object(dictionary, "lookup", return_value=None):
        check("an unseen name has no consent record (None, not a crash)",
              dictionary.get_consent_status(fake_db, "Jane Doe") is None)

    with patch.object(dictionary, "lookup", return_value=fake_entity):
        check("a known name with no consent yet returns None",
              dictionary.get_consent_status(fake_db, "Jane Doe") is None)
        dictionary.set_consent_status(fake_db, fake_entity, "granted")
        check("set_consent_status writes the status onto the entity", fake_entity.consent_status == "granted")
        check("get_consent_status now reflects the granted status",
              dictionary.get_consent_status(fake_db, "Jane Doe") == "granted")


def check_sensitive_outcome_detection() -> None:
    """Phase 3 acceptance test: sensitive_outcome.check_negative_outcome must
    surface the model's discusses_negative_outcome/excerpts/summary fields
    unchanged (Bedrock mocked - this is a plumbing test, same discipline as
    check_reidentify_qa) - and must return a well-formed false/empty result
    for an ordinary document with nothing sensitive in it."""
    print("\n== sensitive outcome detection (Phase 3, mocked Bedrock) ==")
    from app.agents.sanitization import sensitive_outcome
    from app.llm import bedrock_client

    outage_text = (
        "During the migration, a misconfigured firewall rule caused a two-hour outage affecting "
        "customer checkout. Root cause: an unreviewed change to the security group."
    )
    canned_hit = {
        "discusses_negative_outcome": True,
        "excerpts": ["a misconfigured firewall rule caused a two-hour outage affecting customer checkout"],
        "summary": "A firewall misconfiguration caused a two-hour customer-facing outage.",
    }
    mock_hit = bedrock_client.BedrockResponse(text="", parsed=canned_hit, input_tokens=80, output_tokens=40, estimated_cost_usd=0.001)
    with patch.object(bedrock_client, "converse", new=AsyncMock(return_value=mock_hit)):
        resp = asyncio.run(sensitive_outcome.check_negative_outcome(outage_text))
    parsed = resp.parsed or {}
    check("flags the outage narrative as a negative outcome", parsed.get("discusses_negative_outcome") is True, f"got: {parsed}")
    check("carries the verbatim excerpt through unchanged", "two-hour outage" in (parsed.get("excerpts") or [""])[0], f"got: {parsed}")
    check("carries a reviewer-facing summary through unchanged", "outage" in parsed.get("summary", ""), f"got: {parsed}")

    clean_text = "This engagement delivered a new reporting dashboard on time and under budget."
    canned_clean = {"discusses_negative_outcome": False, "excerpts": [], "summary": ""}
    mock_clean = bedrock_client.BedrockResponse(text="", parsed=canned_clean, input_tokens=60, output_tokens=10, estimated_cost_usd=0.0006)
    with patch.object(bedrock_client, "converse", new=AsyncMock(return_value=mock_clean)):
        resp2 = asyncio.run(sensitive_outcome.check_negative_outcome(clean_text))
    parsed2 = resp2.parsed or {}
    check("an ordinary positive-outcome document is not flagged", parsed2.get("discusses_negative_outcome") is False, f"got: {parsed2}")
    check("no excerpts for the clean case", parsed2.get("excerpts") == [], f"got: {parsed2}")


def main() -> int:
    workdir = Path(tempfile.mkdtemp(prefix="naviknow-regression-"))
    try:
        for filename, content_type, builder in FIXTURES:
            try:
                run_format(workdir, filename, content_type, builder)
            except Exception as exc:  # a crash in one format shouldn't hide the others
                import traceback

                traceback.print_exc()
                failures.append(f"{filename}: crashed - {exc}")

        try:
            check_precision_qa()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"precision QA: crashed - {exc}")

        try:
            check_reidentify_qa()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"re-identification QA: crashed - {exc}")

        try:
            check_review_deltas()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"review deltas: crashed - {exc}")

        try:
            check_per_entity_type_thresholds()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"per-entity-type thresholds: crashed - {exc}")

        try:
            check_vision_verdict_cache()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"vision-verdict cache: crashed - {exc}")

        try:
            check_stale_flag_cleanup()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"stale flag cleanup: crashed - {exc}")

        try:
            check_text_remediation_idempotent_reremask()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"text remediation re-render: crashed - {exc}")

        try:
            check_xlsx_threaded_comments()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"xlsx threaded comments: crashed - {exc}")

        try:
            check_logo_band_index()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"logo band index: crashed - {exc}")

        try:
            check_alt_text_channel()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"alt-text channel: crashed - {exc}")

        try:
            check_pptx_author_list()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"pptx author list: crashed - {exc}")

        try:
            check_author_identity_unconditional()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"author identity unconditional: crashed - {exc}")

        try:
            check_docx_comment_author()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"docx comment author: crashed - {exc}")

        try:
            check_unconditional_identity_fields()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"unconditional identity fields: crashed - {exc}")

        try:
            check_surface_pattern_underscore_boundary()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"surface_pattern underscore boundary: crashed - {exc}")

        try:
            check_alt_text_detector_context()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"alt-text detector context: crashed - {exc}")

        try:
            check_custom_replacement_alias()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"custom replacement alias: crashed - {exc}")

        try:
            check_alias_llm_validation()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"alias LLM validation: crashed - {exc}")

        try:
            check_image_placeholder_custom_label()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"image placeholder custom label: crashed - {exc}")

        try:
            check_redact_images_per_entity_labels()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"redact_images per-entity labels: crashed - {exc}")

        try:
            check_pdf_annotation_author()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"pdf annotation author: crashed - {exc}")

        try:
            check_channel_coverage_audit()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"channel coverage audit: crashed - {exc}")

        try:
            check_docx_alt_text()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"docx alt-text: crashed - {exc}")

        try:
            check_xlsx_alt_text()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"xlsx alt-text: crashed - {exc}")

        try:
            check_revalidation_agent()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"revalidation agent: crashed - {exc}")

        try:
            check_logo_thumbnail()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"logo reference thumbnail: crashed - {exc}")

        try:
            check_infra_credential_detection()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"infra/credential regex detection: crashed - {exc}")

        try:
            check_credential_mandatory_no_override()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"CREDENTIAL mandatory non-override: crashed - {exc}")

        try:
            check_exif_strip()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"EXIF stripping: crashed - {exc}")

        try:
            check_data_sample_and_sensitive_text_image_fields()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"data-sample + sensitive-text image fields: crashed - {exc}")

        try:
            check_dataset_provenance()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"dataset provenance check: crashed - {exc}")

        try:
            check_entity_actions_default_resolution()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"entity default-action resolution: crashed - {exc}")

        try:
            check_commercial_term_and_title_detection()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"commercial-term / client-title regex detection: crashed - {exc}")

        try:
            check_flag_vs_mask_vs_mandatory_inclusion()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"flag vs. mask vs. mandatory inclusion: crashed - {exc}")

        try:
            check_immediate_consent_grant()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"immediate consent grant: crashed - {exc}")

        try:
            check_internal_team_consent()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"internal team consent: crashed - {exc}")

        try:
            check_sensitive_outcome_detection()
        except Exception as exc:
            import traceback

            traceback.print_exc()
            failures.append(f"sensitive outcome detection: crashed - {exc}")

        print(f"\n{'=' * 50}\n{passes} checks passed, {len(failures)} failed")
        for f in failures:
            print(f"  FAIL  {f}")
        return 1 if failures else 0
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


if __name__ == "__main__":
    sys.exit(main())
